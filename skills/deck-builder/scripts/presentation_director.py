#!/usr/bin/env python3
"""Local Presentation Director UI for Codex + Presentations workflows.

This script does not generate PPTX files. It creates a small local UI for:
1. intake choices before generation,
2. brief confirmation before calling Presentations,
3. visual revision choices after v1,
4. final version selection,
5. view-only HTML companions from rendered slide previews.
"""

from __future__ import annotations

import argparse
import base64
import hashlib
import html
import json
import mimetypes
import os
import re
import secrets
import shutil
import sys
import threading
import time
import webbrowser
from dataclasses import dataclass, replace
from datetime import datetime
from http import HTTPStatus
from http.server import BaseHTTPRequestHandler, ThreadingHTTPServer
from pathlib import Path
from typing import Any
from urllib.parse import parse_qs, unquote, urlparse


JsonDict = dict[str, Any]

DEFAULT_PORT: int = 8765
DEFAULT_HOST: str = "127.0.0.1"
IMAGE_EXTENSIONS: set[str] = {".png", ".jpg", ".jpeg", ".webp"}
STATUS_FILES: dict[str, str] = {
    "confirmed": "confirmed.ready",
    "revision": "revision.ready",
    "final-selection": "final-selected.ready",
}
PAGE_PATHS: dict[str, str] = {
    "intake": "/intake",
    "visual-inspiration": "/visual-inspiration",
    "confirm": "/confirm",
    "style-review": "/style-review",
    "compare": "/compare",
}
SUPPORTED_UI_LANGUAGES: set[str] = {"zh", "en", "de", "fr", "it", "es"}
HTML_LANG: dict[str, str] = {
    "zh": "zh-CN",
    "en": "en",
    "de": "de",
    "fr": "fr",
    "it": "it",
    "es": "es",
}
UI_COPY: dict[str, dict[str, str]] = {
    "zh": {
        "brief_gate": "Brief Confirmation Gate",
        "confirm_title": "确认生成简报",
        "confirm_intro": "请最后确认一次。只有点击“确认并开始生成”后，agent 才应按所选输出格式开始生成。",
        "topic": "主题",
        "sources": "资料来源",
        "summary": "选择汇总",
        "item": "项目",
        "selection": "选择",
        "source": "来源",
        "output_format": "输出格式",
        "html_transition": "过渡效果",
        "html_animation": "动画密度",
        "html_gradient": "渐变背景",
        "visual_direction": "选定视觉方向",
        "background": "背景策略",
        "layout": "版式节奏",
        "chart": "图表语法",
        "image_strategy": "图片策略",
        "risk": "风险",
        "pre_generation_risks": "生成前风险",
        "generation_strategy": "生成策略",
        "generation_strategy_text": "先生成 v1 PPTX 和 contact sheet，并集中保存到 {task_dir}，然后打开 style-review.html 供选择是否重绘。",
        "back_visual": "返回修改视觉方向",
        "confirm_button": "确认并开始生成",
        "no_risks": "未发现明显风险。",
        "no_sources": "未记录资料来源。可以在下面粘贴本地路径、网页 URL 或 Google Drive 地址。",
        "confirmed_title": "Brief confirmed",
        "confirmed_message": "确认已收到。你不需要回到聊天里回复；agent 会检测 confirmed.ready 并自动继续生成。",
        "invalid_token": "Missing or invalid confirmation token. Open the confirmation page and submit the form.",
        "default": "default",
        "user_selected": "user-selected",
        "unknown": "unknown",
    },
    "en": {
        "brief_gate": "Brief Confirmation Gate",
        "confirm_title": "Confirm Generation Brief",
        "confirm_intro": "Please review the plan one last time. The agent should start generation for the selected output format only after you click \"Confirm and start generation.\"",
        "topic": "Topic",
        "sources": "Source Material",
        "summary": "Selection Summary",
        "item": "Item",
        "selection": "Selection",
        "source": "Source",
        "output_format": "Output Format",
        "html_transition": "Transition",
        "html_animation": "Animation density",
        "html_gradient": "Gradient background",
        "visual_direction": "Selected Visual Direction",
        "background": "Background Strategy",
        "layout": "Layout Rhythm",
        "chart": "Chart Grammar",
        "image_strategy": "Image Strategy",
        "risk": "Risk",
        "pre_generation_risks": "Pre-Generation Risks",
        "generation_strategy": "Generation Strategy",
        "generation_strategy_text": "Generate the v1 PPTX and contact sheet first, save them under {task_dir}, then open style-review.html so you can decide whether to redraw the deck.",
        "back_visual": "Back to visual direction",
        "confirm_button": "Confirm and start generation",
        "no_risks": "No obvious risks detected.",
        "no_sources": "No source material has been recorded. Add local paths, web URLs, or Google Drive links before generation.",
        "confirmed_title": "Brief confirmed",
        "confirmed_message": "Confirmed. You do not need to reply in chat; the agent will detect confirmed.ready and continue automatically.",
        "invalid_token": "Missing or invalid confirmation token. Open the confirmation page and submit the form.",
        "default": "default",
        "user_selected": "user-selected",
        "unknown": "unknown",
    },
    "de": {
        "brief_gate": "Bestätigung des Briefings",
        "confirm_title": "Generierungsbrief bestätigen",
        "confirm_intro": "Bitte prüfen Sie den Plan ein letztes Mal. Erst nach dem Klick auf \"Bestätigen und Generierung starten\" sollte der Agent die Generierung im gewählten Ausgabeformat starten.",
        "topic": "Thema",
        "sources": "Quellenmaterial",
        "summary": "Zusammenfassung der Auswahl",
        "item": "Punkt",
        "selection": "Auswahl",
        "source": "Quelle",
        "output_format": "Ausgabeformat",
        "html_transition": "Übergang",
        "html_animation": "Animationsdichte",
        "html_gradient": "Verlaufshintergrund",
        "visual_direction": "Ausgewählte visuelle Richtung",
        "background": "Hintergrundstrategie",
        "layout": "Layoutrhythmus",
        "chart": "Diagrammregeln",
        "image_strategy": "Bildstrategie",
        "risk": "Risiko",
        "pre_generation_risks": "Risiken vor der Generierung",
        "generation_strategy": "Generierungsstrategie",
        "generation_strategy_text": "Zuerst werden v1-PPTX und Contact Sheet erzeugt und unter {task_dir} gespeichert. Danach wird style-review.html geöffnet, damit Sie entscheiden können, ob das Deck visuell überarbeitet werden soll.",
        "back_visual": "Zur visuellen Richtung zurück",
        "confirm_button": "Bestätigen und Generierung starten",
        "no_risks": "Keine offensichtlichen Risiken erkannt.",
        "no_sources": "Es wurden keine Quellen erfasst. Fügen Sie vor der Generierung lokale Pfade, Web-URLs oder Google-Drive-Links hinzu.",
        "confirmed_title": "Briefing bestätigt",
        "confirmed_message": "Bestätigt. Sie müssen nicht im Chat antworten; der Agent erkennt confirmed.ready und fährt automatisch fort.",
        "invalid_token": "Fehlendes oder ungültiges Bestätigungstoken. Öffnen Sie die Bestätigungsseite und senden Sie das Formular ab.",
        "default": "Standard",
        "user_selected": "vom Benutzer gewählt",
        "unknown": "unbekannt",
    },
    "fr": {
        "brief_gate": "Validation du brief",
        "confirm_title": "Confirmer le brief de génération",
        "confirm_intro": "Veuillez relire le plan une dernière fois. L'agent ne doit lancer la génération du format choisi qu'après votre clic sur \"Confirmer et lancer la génération\".",
        "topic": "Sujet",
        "sources": "Sources",
        "summary": "Résumé des choix",
        "item": "Élément",
        "selection": "Choix",
        "source": "Source",
        "output_format": "Format de sortie",
        "html_transition": "Transition",
        "html_animation": "Densité d'animation",
        "html_gradient": "Fond dégradé",
        "visual_direction": "Direction visuelle sélectionnée",
        "background": "Stratégie de fond",
        "layout": "Rythme de mise en page",
        "chart": "Grammaire des graphiques",
        "image_strategy": "Stratégie d'image",
        "risk": "Risque",
        "pre_generation_risks": "Risques avant génération",
        "generation_strategy": "Stratégie de génération",
        "generation_strategy_text": "Générer d'abord le PPTX v1 et la planche de contact, les enregistrer dans {task_dir}, puis ouvrir style-review.html pour décider d'une éventuelle refonte visuelle.",
        "back_visual": "Retour à la direction visuelle",
        "confirm_button": "Confirmer et lancer la génération",
        "no_risks": "Aucun risque évident détecté.",
        "no_sources": "Aucune source n'a été enregistrée. Ajoutez des chemins locaux, des URL web ou des liens Google Drive avant la génération.",
        "confirmed_title": "Brief confirmé",
        "confirmed_message": "Confirmé. Vous n'avez pas besoin de répondre dans le chat; l'agent détectera confirmed.ready et continuera automatiquement.",
        "invalid_token": "Jeton de confirmation manquant ou invalide. Ouvrez la page de confirmation et envoyez le formulaire.",
        "default": "par défaut",
        "user_selected": "choisi par l'utilisateur",
        "unknown": "inconnu",
    },
    "it": {
        "brief_gate": "Conferma del brief",
        "confirm_title": "Conferma il brief di generazione",
        "confirm_intro": "Rivedi il piano un'ultima volta. L'agente dovrebbe avviare la generazione del formato scelto solo dopo il clic su \"Conferma e avvia la generazione\".",
        "topic": "Argomento",
        "sources": "Fonti",
        "summary": "Riepilogo delle scelte",
        "item": "Voce",
        "selection": "Scelta",
        "source": "Fonte",
        "output_format": "Formato di output",
        "html_transition": "Transizione",
        "html_animation": "Densità animazione",
        "html_gradient": "Sfondo sfumato",
        "visual_direction": "Direzione visiva selezionata",
        "background": "Strategia di sfondo",
        "layout": "Ritmo del layout",
        "chart": "Grammatica dei grafici",
        "image_strategy": "Strategia immagini",
        "risk": "Rischio",
        "pre_generation_risks": "Rischi prima della generazione",
        "generation_strategy": "Strategia di generazione",
        "generation_strategy_text": "Genera prima il PPTX v1 e il contact sheet, salvali in {task_dir}, poi apri style-review.html per decidere se ridisegnare il deck.",
        "back_visual": "Torna alla direzione visiva",
        "confirm_button": "Conferma e avvia la generazione",
        "no_risks": "Nessun rischio evidente rilevato.",
        "no_sources": "Nessuna fonte registrata. Aggiungi percorsi locali, URL web o link Google Drive prima della generazione.",
        "confirmed_title": "Brief confermato",
        "confirmed_message": "Confermato. Non serve rispondere in chat; l'agente rileverà confirmed.ready e continuerà automaticamente.",
        "invalid_token": "Token di conferma mancante o non valido. Apri la pagina di conferma e invia il modulo.",
        "default": "predefinito",
        "user_selected": "scelto dall'utente",
        "unknown": "sconosciuto",
    },
    "es": {
        "brief_gate": "Confirmación del brief",
        "confirm_title": "Confirmar el brief de generación",
        "confirm_intro": "Revisa el plan una última vez. El agente solo debe iniciar la generación del formato elegido después de que hagas clic en \"Confirmar e iniciar generación\".",
        "topic": "Tema",
        "sources": "Fuentes",
        "summary": "Resumen de selecciones",
        "item": "Elemento",
        "selection": "Selección",
        "source": "Fuente",
        "output_format": "Formato de salida",
        "html_transition": "Transición",
        "html_animation": "Densidad de animación",
        "html_gradient": "Fondo degradado",
        "visual_direction": "Dirección visual seleccionada",
        "background": "Estrategia de fondo",
        "layout": "Ritmo de diseño",
        "chart": "Gramática de gráficos",
        "image_strategy": "Estrategia de imágenes",
        "risk": "Riesgo",
        "pre_generation_risks": "Riesgos antes de generar",
        "generation_strategy": "Estrategia de generación",
        "generation_strategy_text": "Primero genera el PPTX v1 y la hoja de contacto, guárdalos en {task_dir}, y luego abre style-review.html para decidir si redibujar el deck.",
        "back_visual": "Volver a dirección visual",
        "confirm_button": "Confirmar e iniciar generación",
        "no_risks": "No se detectaron riesgos evidentes.",
        "no_sources": "No se registraron fuentes. Añade rutas locales, URL web o enlaces de Google Drive antes de generar.",
        "confirmed_title": "Brief confirmado",
        "confirmed_message": "Confirmado. No necesitas responder en el chat; el agente detectará confirmed.ready y continuará automáticamente.",
        "invalid_token": "Token de confirmación ausente o no válido. Abre la página de confirmación y envía el formulario.",
        "default": "predeterminado",
        "user_selected": "seleccionado por el usuario",
        "unknown": "desconocido",
    },
}
ADDITIONAL_UI_COPY: dict[str, dict[str, str]] = {
    "zh": {
        "brief_gate": "简报确认门禁",
        "confirmed_title": "简报已确认",
        "invalid_token": "确认令牌缺失或无效。请打开确认页并提交表单。",
        "default": "默认",
        "user_selected": "用户选择",
        "no_sources": "未记录资料来源。可以在下面填写本地路径、网页 URL 或 Google Drive 地址。",
        "intake_topline": "Presentation Director",
        "intake_title": "生成前信息收集",
        "intake_intro": "先确认会影响 PPTX 质量的关键信息。每题都有默认推荐，也可以选择自定义。",
        "source_material": "资料来源",
        "source_paths_label": "资料路径 / 网页 / Google Drive 地址",
        "source_paths_placeholder": "每行一个来源。例如：\n/Users/you/project/docs\nhttps://example.com/report\nhttps://drive.google.com/file/d/...\nhttps://docs.google.com/document/d/...",
        "source_paths_meta": "可以填写本地文件夹、本地文件、普通网页 URL、Google Drive / Docs / Slides / Sheets 地址。Google Drive 地址会作为来源链接记录，后续由 agent 按权限读取或要求你授权。",
        "topic_title_label": "主题 / 标题",
        "extra_notes": "额外说明",
        "extra_notes_placeholder": "例如必须保留的页面、禁用内容、特殊听众背景。",
        "custom_placeholder": "输入你的自定义说明",
        "next_visual": "下一步：视觉候选",
        "visual_gate": "视觉方向门禁",
        "visual_title": "选择第一版视觉方向",
        "visual_intro": "这些候选会根据主题、PPT 类型和听众动态生成。它们借鉴 design-lock、ui-ux-pro-max、HTML deck/theme catalog 的做法，但最终仍交给 Codex Presentations 生成可编辑 PPTX。",
        "current_topic": "当前主题",
        "visual_notes": "视觉补充要求",
        "visual_notes_placeholder": "例如：更像顶级咨询公司、更少卡片、背景更有层次、适合医学研究听众。",
        "back_intake": "返回修改 intake",
        "next_confirm": "下一步：汇总确认",
        "best_for": "适合",
        "inspiration": "借鉴",
        "evidence_page": "证据页",
        "style_review": "视觉复审",
        "style_title": "视觉复审",
        "style_intro": "基于当前最新版本 {version_name} 的 contact sheet 选择是否生成对比版本。不要复制 JSON，点击按钮即可。",
        "current_version": "当前版本",
        "missing_contact_sheet": "还没有找到 {path}。生成版本后刷新此页。",
        "missing_qa_summary": "暂无 QA 摘要。",
        "revision_count_title": "生成几个对比版本?",
        "keep_current_version": "保持 v1，进入最终选择",
        "one_revision": "生成一个对比版本",
        "two_revisions": "生成两个对比版本",
        "revision_notes_placeholder": "补充说明，例如：第 5 页架构图需要更清楚。",
        "confirm_visual_choice": "确认视觉选择",
        "version_compare": "版本比较",
        "compare_title": "选择最终版本",
        "choose_after_action": "选择后动作",
        "final_notes_placeholder": "最终选择理由或仍需注意的问题。",
        "confirm_final_version": "确认最终版本",
        "continue_editing": "继续修改",
        "choose_version": "选择 {version}",
        "no_contact_sheet": "没有 contact sheet。",
        "no_versions": "还没有可比较版本。请先生成 v1。",
        "pptx_label": "PPTX",
        "revision_saved_title": "修改选择已保存",
        "revision_saved_message": "修改选择已收到。你不需要回到聊天里回复；agent 会检测 revision.ready 并自动生成对比版本。",
        "final_selected_title": "最终版本已选择",
        "final_selected_message": "最终版本选择已收到。你不需要回到聊天里回复；agent 会检测 final-selected.ready 并自动做最终交付。",
        "nav_intake": "信息收集",
        "nav_visual": "视觉方向",
        "nav_style": "视觉复审",
        "nav_compare": "版本比较",
    },
    "en": {
        "intake_topline": "Presentation Director",
        "intake_title": "Pre-Generation Intake",
        "intake_intro": "Confirm the key details that affect PPTX quality. Each question has a recommended default, and you can choose a custom answer.",
        "source_material": "Source Material",
        "source_paths_label": "Source paths / web pages / Google Drive links",
        "source_paths_placeholder": "One source per line. For example:\n/Users/you/project/docs\nhttps://example.com/report\nhttps://drive.google.com/file/d/...\nhttps://docs.google.com/document/d/...",
        "source_paths_meta": "You can enter local folders, local files, regular web URLs, and Google Drive / Docs / Slides / Sheets links. Google Drive links are recorded as source links; the agent will read them according to permissions or ask for authorization.",
        "topic_title_label": "Topic / title",
        "extra_notes": "Additional Notes",
        "extra_notes_placeholder": "For example: pages to preserve, forbidden content, or special audience context.",
        "custom_placeholder": "Enter your custom note",
        "next_visual": "Next: visual candidates",
        "visual_gate": "Visual Inspiration Gate",
        "visual_title": "Choose the First-Draft Visual Direction",
        "visual_intro": "These candidates are generated from the topic, PPT type, and audience. They borrow from design-locks, ui-ux-pro-max, and HTML deck/theme catalog patterns, while Codex Presentations still creates the editable PPTX.",
        "current_topic": "Current Topic",
        "visual_notes": "Additional Visual Requirements",
        "visual_notes_placeholder": "For example: more like a top-tier consulting deck, fewer cards, richer backgrounds, or suitable for a medical research audience.",
        "back_intake": "Back to intake",
        "next_confirm": "Next: summary confirmation",
        "best_for": "Best for",
        "inspiration": "Inspiration",
        "evidence_page": "Evidence page",
        "style_review": "Style Review",
        "style_title": "Style Review",
        "style_intro": "Use the contact sheet for the latest version {version_name} to decide whether to generate comparison versions. No JSON copying is needed; just click the button.",
        "current_version": "Current Version",
        "missing_contact_sheet": "Could not find {path}. Refresh this page after the version is generated.",
        "missing_qa_summary": "No QA summary yet.",
        "revision_count_title": "How many comparison versions should be generated?",
        "keep_current_version": "Keep v1 and move to final selection",
        "one_revision": "Generate one comparison version",
        "two_revisions": "Generate two comparison versions",
        "revision_notes_placeholder": "Add notes, for example: make the architecture diagram on slide 5 clearer.",
        "confirm_visual_choice": "Confirm visual choice",
        "version_compare": "Version Compare",
        "compare_title": "Choose the Final Version",
        "choose_after_action": "After Selection",
        "final_notes_placeholder": "Reason for the final choice or remaining issues to watch.",
        "confirm_final_version": "Confirm final version",
        "continue_editing": "Continue editing",
        "choose_version": "Choose {version}",
        "no_contact_sheet": "No contact sheet found.",
        "no_versions": "No comparable versions yet. Generate v1 first.",
        "pptx_label": "PPTX",
        "revision_saved_title": "Revision saved",
        "revision_saved_message": "Revision choices received. You do not need to reply in chat; the agent will detect revision.ready and generate comparison versions automatically.",
        "final_selected_title": "Final version selected",
        "final_selected_message": "Final version selection received. You do not need to reply in chat; the agent will detect final-selected.ready and complete the final delivery automatically.",
        "nav_intake": "Intake",
        "nav_visual": "Visual Direction",
        "nav_style": "Style Review",
        "nav_compare": "Compare",
    },
    "de": {
        "intake_topline": "Presentation Director",
        "intake_title": "Informationen vor der Generierung",
        "intake_intro": "Bestätigen Sie zuerst die wichtigsten Angaben, die die PPTX-Qualität beeinflussen. Jede Frage hat eine empfohlene Standardeinstellung; Sie können auch eine eigene Antwort wählen.",
        "source_material": "Quellenmaterial",
        "source_paths_label": "Quellenpfade / Webseiten / Google-Drive-Links",
        "source_paths_placeholder": "Eine Quelle pro Zeile. Zum Beispiel:\n/Users/you/project/docs\nhttps://example.com/report\nhttps://drive.google.com/file/d/...\nhttps://docs.google.com/document/d/...",
        "source_paths_meta": "Sie können lokale Ordner, lokale Dateien, normale Web-URLs und Google Drive / Docs / Slides / Sheets-Links eintragen. Google-Drive-Links werden als Quellenlinks gespeichert; der Agent liest sie je nach Berechtigung oder fordert eine Autorisierung an.",
        "topic_title_label": "Thema / Titel",
        "extra_notes": "Zusätzliche Hinweise",
        "extra_notes_placeholder": "Zum Beispiel: Seiten, die erhalten bleiben müssen, verbotene Inhalte oder besonderer Kontext zur Zielgruppe.",
        "custom_placeholder": "Eigene Anmerkung eingeben",
        "next_visual": "Weiter: visuelle Kandidaten",
        "visual_gate": "Tor für visuelle Richtung",
        "visual_title": "Visuelle Richtung für den ersten Entwurf wählen",
        "visual_intro": "Diese Kandidaten werden aus Thema, PPT-Typ und Zielgruppe abgeleitet. Sie nutzen Muster aus design-locks, ui-ux-pro-max und HTML deck/theme catalogs; Codex Presentations erstellt daraus weiterhin eine editierbare PPTX.",
        "current_topic": "Aktuelles Thema",
        "visual_notes": "Zusätzliche visuelle Anforderungen",
        "visual_notes_placeholder": "Zum Beispiel: näher an einer Top-Consulting-Präsentation, weniger Karten, mehr Tiefe im Hintergrund oder passend für ein medizinisches Forschungspublikum.",
        "back_intake": "Zurück zur Informationsabfrage",
        "next_confirm": "Weiter: Zusammenfassung bestätigen",
        "best_for": "Geeignet für",
        "inspiration": "Inspiration",
        "evidence_page": "Evidenzseite",
        "style_review": "Visuelle Prüfung",
        "style_title": "Visuelle Prüfung",
        "style_intro": "Nutzen Sie das Contact Sheet der neuesten Version {version_name}, um zu entscheiden, ob Vergleichsversionen erzeugt werden sollen. Kein JSON-Kopieren nötig; klicken Sie einfach auf die Schaltfläche.",
        "current_version": "Aktuelle Version",
        "missing_contact_sheet": "{path} wurde noch nicht gefunden. Aktualisieren Sie diese Seite nach der Versionserstellung.",
        "missing_qa_summary": "Noch keine QA-Zusammenfassung vorhanden.",
        "revision_count_title": "Wie viele Vergleichsversionen sollen erzeugt werden?",
        "keep_current_version": "v1 behalten und zur finalen Auswahl gehen",
        "one_revision": "Eine Vergleichsversion erzeugen",
        "two_revisions": "Zwei Vergleichsversionen erzeugen",
        "revision_notes_placeholder": "Zusätzliche Hinweise, zum Beispiel: Das Architekturdiagramm auf Folie 5 soll klarer werden.",
        "confirm_visual_choice": "Visuelle Auswahl bestätigen",
        "version_compare": "Versionsvergleich",
        "compare_title": "Finale Version wählen",
        "choose_after_action": "Aktion nach der Auswahl",
        "final_notes_placeholder": "Begründung der finalen Auswahl oder verbleibende Punkte.",
        "confirm_final_version": "Finale Version bestätigen",
        "continue_editing": "Weiter bearbeiten",
        "choose_version": "{version} wählen",
        "no_contact_sheet": "Kein Contact Sheet gefunden.",
        "no_versions": "Es gibt noch keine vergleichbaren Versionen. Erzeugen Sie zuerst v1.",
        "pptx_label": "PPTX",
        "revision_saved_title": "Überarbeitung gespeichert",
        "revision_saved_message": "Die Überarbeitungsauswahl wurde empfangen. Sie müssen nicht im Chat antworten; der Agent erkennt revision.ready und erzeugt automatisch Vergleichsversionen.",
        "final_selected_title": "Finale Version ausgewählt",
        "final_selected_message": "Die finale Versionsauswahl wurde empfangen. Sie müssen nicht im Chat antworten; der Agent erkennt final-selected.ready und erstellt automatisch die finale Lieferung.",
        "nav_intake": "Informationsabfrage",
        "nav_visual": "Visuelle Richtung",
        "nav_style": "Visuelle Prüfung",
        "nav_compare": "Vergleich",
    },
}

for language, copy_items in ADDITIONAL_UI_COPY.items():
    UI_COPY.setdefault(language, {}).update(copy_items)
QUESTION_TITLE_L10N: dict[str, dict[str, str]] = {
    "en": {
        "deck_type": "PPT Type",
        "output_format": "Output Format",
        "research_strategy": "Research Strategy",
        "audience": "Audience",
        "goal": "Goal",
        "source_boundary": "Source Boundary",
        "content_language": "Content Language",
        "output_constraints": "Output Constraints",
        "logo_policy": "Logo / Brand Assets",
        "image_policy": "AI Image Policy",
        "visual_freedom": "First-Draft Visual Direction",
        "reference_deck": "Reference Deck",
    },
    "de": {
        "deck_type": "PPT-Typ",
        "output_format": "Ausgabeformat",
        "research_strategy": "Recherche-Strategie",
        "audience": "Zielgruppe",
        "goal": "Ziel",
        "source_boundary": "Quellengrenzen",
        "content_language": "Inhaltssprache",
        "output_constraints": "Umfang und Dauer",
        "logo_policy": "Logo / Markenmaterial",
        "image_policy": "KI-Bildrichtlinie",
        "visual_freedom": "Visuelle Richtung des ersten Entwurfs",
        "reference_deck": "Referenzdeck",
    },
    "fr": {
        "deck_type": "Type de PPT",
        "output_format": "Format de sortie",
        "research_strategy": "Stratégie de recherche",
        "audience": "Public",
        "goal": "Objectif",
        "source_boundary": "Limites des sources",
        "content_language": "Langue du contenu",
        "output_constraints": "Contraintes de sortie",
        "logo_policy": "Logo / actifs de marque",
        "image_policy": "Politique d'images IA",
        "visual_freedom": "Direction visuelle du premier jet",
        "reference_deck": "Deck de référence",
    },
    "it": {
        "deck_type": "Tipo di PPT",
        "output_format": "Formato di output",
        "research_strategy": "Strategia di ricerca",
        "audience": "Pubblico",
        "goal": "Obiettivo",
        "source_boundary": "Limiti delle fonti",
        "content_language": "Lingua dei contenuti",
        "output_constraints": "Vincoli di output",
        "logo_policy": "Logo / asset del brand",
        "image_policy": "Policy immagini IA",
        "visual_freedom": "Direzione visiva della prima bozza",
        "reference_deck": "Deck di riferimento",
    },
    "es": {
        "deck_type": "Tipo de PPT",
        "output_format": "Formato de salida",
        "research_strategy": "Estrategia de investigación",
        "audience": "Audiencia",
        "goal": "Objetivo",
        "source_boundary": "Límites de fuentes",
        "content_language": "Idioma del contenido",
        "output_constraints": "Restricciones de salida",
        "logo_policy": "Logo / activos de marca",
        "image_policy": "Política de imágenes IA",
        "visual_freedom": "Dirección visual del primer borrador",
        "reference_deck": "Deck de referencia",
    },
}
CHOICE_LABEL_L10N: dict[str, dict[str, dict[str, str]]] = {
    "en": {
        "deck_type": {
            "project-report": "Project update",
            "engineering-platform": "Engineering / technical solution",
            "investor-pitch": "Investor / pitch deck",
            "knowledge-teaching": "Academic / course / knowledge explanation",
            "sales-product": "Sales / product presentation",
            "custom": "Custom",
        },
        "output_format": {
            "html-revealjs": "HTML (Reveal.js)",
            "pptx": "PPTX (PowerPoint)",
            "both": "HTML + PPTX (Both)",
        },
        "research_strategy": {
            "hybrid-deep-research": "Hybrid: external Deep Research + Codex verification",
            "codex-web-deep": "Codex deep web research",
            "external-deep-research": "Gemini / Perplexity Deep Research packet",
            "provided-materials": "Use only the materials I provide",
            "custom": "Custom",
        },
        "audience": {
            "executives": "Executives / decision-makers",
            "investors-reviewers": "Investors / reviewers / pitch audience",
            "technical-leaders": "Technical team / engineering reviewers",
            "customers-sales": "Customers / sales audience",
            "teachers-researchers": "Students / teachers / researchers",
            "custom": "Custom",
        },
        "goal": {
            "understand-topic": "Help the audience understand a topic quickly",
            "decision": "Persuade the audience to make a decision",
            "progress-risk": "Report progress, results, and risks",
            "teaching": "Teach / explain knowledge",
            "explain-value": "Show product / project value",
            "custom": "Custom",
        },
        "source_boundary": {
            "provided-only": "Strictly use my provided materials",
            "web-with-sources": "May supplement from the web with cited sources",
            "existing-doc": "Use an existing PPT / document as the content base",
            "reference-quality": "Use a reference deck as the quality and style bar",
            "custom": "Custom",
        },
        "content_language": {
            "zh": "Chinese",
            "en": "English",
            "de": "German",
            "fr": "French",
            "it": "Italian",
            "es": "Spanish",
            "bilingual": "Bilingual",
            "custom": "Custom",
        },
        "output_constraints": {
            "pages-8-10": "8-10 slides",
            "pages-10-12": "10-12 slides",
            "pages-15-20": "15-20 slides",
            "custom": "Custom duration and slide count",
        },
        "logo_policy": {
            "none": "Do not use logos",
            "provided-only": "Only use logos / images I provide",
            "official-sources": "May find official logos and assets",
            "cover-final-only": "Use logos only on cover and final slides",
            "custom": "Custom",
        },
        "image_policy": {
            "none": "Do not use AI-generated images",
            "abstract-only": "Allow abstract backgrounds / concept images only",
            "cover-section": "Allow cover and section images",
            "ask-before-use": "Ask me before each generated image",
            "custom": "Custom",
        },
        "visual_freedom": {
            "delegate": "AI-driven — let the generation engine choose freely",
            "restrained": "More formal and restrained",
            "technical": "More technical / engineering-oriented",
            "investor": "More investor-pitch / high-contrast",
            "academic-editorial": "More academic / editorial",
            "custom": "Custom",
        },
        "reference_deck": {
            "none": "No reference; generate from content",
            "quality-only": "Reference exists, but only as a quality bar",
            "visual-style": "Reference exists; get close to its visual style",
            "existing-ppt": "Existing PPT needs to be revised",
            "custom": "Custom",
        },
    },
    "de": {
        "deck_type": {
            "project-report": "Projektbericht",
            "engineering-platform": "Technische Lösung / Architektur",
            "investor-pitch": "Investor- / Pitch-Deck",
            "knowledge-teaching": "Akademische / Kurs- / Wissensvermittlung",
            "sales-product": "Vertrieb / Produktpräsentation",
            "custom": "Benutzerdefiniert",
        },
        "output_format": {
            "html-revealjs": "HTML (Reveal.js)",
            "pptx": "PPTX (PowerPoint)",
            "both": "HTML + PPTX (Beide)",
        },
        "research_strategy": {
            "hybrid-deep-research": "Hybrid: externe Deep Research + Codex-Prüfung",
            "codex-web-deep": "Codex Deep-Web-Recherche",
            "external-deep-research": "Gemini / Perplexity Deep-Research-Paket",
            "provided-materials": "Nur bereitgestellte Materialien verwenden",
            "custom": "Benutzerdefiniert",
        },
        "audience": {
            "executives": "Führungskräfte / Entscheider",
            "investors-reviewers": "Investoren / Gutachter / Pitch-Publikum",
            "technical-leaders": "Technisches Team / Engineering Review",
            "customers-sales": "Kunden / Vertriebspublikum",
            "teachers-researchers": "Studierende / Lehrende / Forschende",
            "custom": "Benutzerdefiniert",
        },
        "goal": {
            "understand-topic": "Ein Thema schnell verständlich machen",
            "decision": "Zu einer Entscheidung überzeugen",
            "progress-risk": "Fortschritt, Ergebnisse und Risiken berichten",
            "teaching": "Lehren / Wissen vermitteln",
            "explain-value": "Produkt- / Projektwert zeigen",
            "custom": "Benutzerdefiniert",
        },
        "source_boundary": {
            "provided-only": "Ausschließlich bereitgestellte Materialien verwenden",
            "web-with-sources": "Web-Ergänzungen mit Quellenangaben erlaubt",
            "existing-doc": "Bestehendes PPT / Dokument als Inhaltsbasis nutzen",
            "reference-quality": "Referenzdeck als Qualitäts- und Stilmaßstab nutzen",
            "custom": "Benutzerdefiniert",
        },
        "content_language": {
            "zh": "Chinesisch",
            "en": "Englisch",
            "de": "Deutsch",
            "fr": "Französisch",
            "it": "Italienisch",
            "es": "Spanisch",
            "bilingual": "Zweisprachig",
            "custom": "Benutzerdefiniert",
        },
        "output_constraints": {
            "pages-8-10": "8-10 Folien",
            "pages-10-12": "10-12 Folien",
            "pages-15-20": "15-20 Folien",
            "custom": "Benutzerdefinierte Dauer und Folienzahl",
        },
        "logo_policy": {
            "none": "Keine Logos verwenden",
            "provided-only": "Nur bereitgestellte Logos / Bilder verwenden",
            "official-sources": "Offizielle Logos und Assets dürfen gesucht werden",
            "cover-final-only": "Logos nur auf Titelfolie und Schlussfolie",
            "custom": "Benutzerdefiniert",
        },
        "image_policy": {
            "none": "Keine KI-generierten Bilder verwenden",
            "abstract-only": "Nur abstrakte Hintergründe / Konzeptbilder erlauben",
            "cover-section": "Titel- und Kapitelbilder erlauben",
            "ask-before-use": "Vor jedem generierten Bild fragen",
            "custom": "Benutzerdefiniert",
        },
        "visual_freedom": {
            "delegate": "KI-gesteuert — Generierungs-Engine frei wählen lassen",
            "restrained": "Formeller und zurückhaltender",
            "technical": "Technischer / stärker engineering-orientiert",
            "investor": "Mehr Investor-Pitch / hoher Kontrast",
            "academic-editorial": "Akademischer / editorialer",
            "custom": "Benutzerdefiniert",
        },
        "reference_deck": {
            "none": "Keine Referenz; aus dem Inhalt generieren",
            "quality-only": "Referenz vorhanden, nur als Qualitätsmaßstab",
            "visual-style": "Referenz vorhanden; visuellen Stil annähern",
            "existing-ppt": "Bestehendes PPT soll überarbeitet werden",
            "custom": "Benutzerdefiniert",
        },
    },
}
ADDITIONAL_QUESTION_TITLE_L10N: dict[str, dict[str, str]] = {
    "en": {
        "palette_direction": "Palette Direction",
        "structure_direction": "Structure Rhythm",
        "visual_expression": "Visual Expression",
        "image_strategy": "Image Strategy",
        "logo_strategy": "Logo / Brand Presence",
    },
    "de": {
        "palette_direction": "Farbpalette",
        "structure_direction": "Strukturrhythmus",
        "visual_expression": "Visueller Ausdruck",
        "image_strategy": "Bildstrategie",
        "logo_strategy": "Logo / Markenpräsenz",
    },
}

QUESTION_PROMPT_L10N: dict[str, dict[str, str]] = {
    "en": {
        "deck_type": "What kind of PPT are you creating?",
        "output_format": "Which presentation output format should be generated?",
        "research_strategy": "If the material is incomplete, how should research material be gathered first?",
        "audience": "Who is the main audience for this PPT?",
        "goal": "What is the main goal of this PPT?",
        "source_boundary": "How should source material be used?",
        "content_language": "What language should the slide body use?",
        "output_constraints": "What are the slide count and presentation duration constraints?",
        "logo_policy": "How should logos and brand assets be handled?",
        "image_policy": "Should AI-generated images be allowed?",
        "visual_freedom": "How should the first-draft visual direction be handled?",
        "reference_deck": "Do you have a reference deck or style sample?",
        "palette_direction": "Should the palette direction change?",
        "structure_direction": "Should the structure rhythm change?",
        "visual_expression": "Should the visual expression change?",
        "image_strategy": "How should images and visual assets be adjusted?",
        "logo_strategy": "How should logos and brand presence be adjusted?",
    },
    "de": {
        "deck_type": "Welche Art von PPT soll erstellt werden?",
        "output_format": "Welches Präsentationsformat soll erzeugt werden?",
        "research_strategy": "Wenn das Material unvollständig ist, wie sollen zuerst Recherchematerialien beschafft werden?",
        "audience": "Für wen ist diese PPT hauptsächlich gedacht?",
        "goal": "Was ist das Hauptziel dieser PPT?",
        "source_boundary": "Wie soll das Quellenmaterial verwendet werden?",
        "content_language": "Welche Sprache soll der Folientext verwenden?",
        "output_constraints": "Welche Vorgaben gibt es für Folienzahl und Vortragsdauer?",
        "logo_policy": "Wie sollen Logos und Markenmaterial behandelt werden?",
        "image_policy": "Sollen KI-generierte Bilder erlaubt sein?",
        "visual_freedom": "Wie soll die visuelle Richtung des ersten Entwurfs behandelt werden?",
        "reference_deck": "Gibt es ein Referenzdeck oder Stilbeispiele?",
        "palette_direction": "Soll sich die Farbpalette ändern?",
        "structure_direction": "Soll sich der Strukturrhythmus ändern?",
        "visual_expression": "Soll sich der visuelle Ausdruck ändern?",
        "image_strategy": "Wie sollen Bilder und visuelle Materialien angepasst werden?",
        "logo_strategy": "Wie sollen Logos und Markenpräsenz angepasst werden?",
    },
}

ADDITIONAL_CHOICE_LABEL_L10N: dict[str, dict[str, dict[str, str]]] = {
    "en": {
        "palette_direction": {
            "keep": "Keep current",
            "warm-editorial": "Warm paper-like / editorial",
            "engineering-blue-gray": "Engineering blue-gray / technical",
            "high-contrast-pitch": "High-contrast / pitch style",
            "brand-boost": "Strengthen brand colors",
            "custom": "Custom",
        },
        "structure_direction": {
            "keep": "Keep current",
            "stronger-story": "Stronger story rhythm",
            "more-analytical": "More analytical",
            "more-architecture": "More architecture / system diagrams",
            "more-product-demo": "More product-launch / demo feel",
            "denser-internal": "Denser internal review style",
            "custom": "Custom",
        },
        "visual_expression": {
            "keep": "Keep current",
            "more-premium": "More premium / more designed",
            "more-restrained": "More restrained / less decorative",
            "less-cards": "Fewer cards, more open composition",
            "stronger-evidence": "Stronger evidence-led storytelling",
            "custom": "Custom",
        },
        "image_strategy": {
            "keep": "Keep current",
            "fewer-decorative": "Use fewer decorative images",
            "more-ai-concept": "Add AI concept imagery",
            "official-only": "Use only real screenshots / official assets",
            "stronger-cover-section": "Strengthen cover or section visuals",
            "custom": "Custom",
        },
        "logo_strategy": {
            "keep": "Keep current",
            "cover-final-only": "Keep logos only on cover and final slides",
            "footer-brand": "Strengthen footer branding",
            "partner-logo-page": "Add partner / school / company logo page",
            "remove-all": "Remove all logos",
            "custom": "Custom",
        },
    },
    "de": {
        "palette_direction": {
            "keep": "Aktuell beibehalten",
            "warm-editorial": "Warmer Papierlook / Editorial",
            "engineering-blue-gray": "Engineering-Blaugrau / technisch",
            "high-contrast-pitch": "Hoher Kontrast / Pitch-Stil",
            "brand-boost": "Markenfarben stärken",
            "custom": "Benutzerdefiniert",
        },
        "structure_direction": {
            "keep": "Aktuell beibehalten",
            "stronger-story": "Stärkerer Story-Rhythmus",
            "more-analytical": "Analytischer",
            "more-architecture": "Mehr Architektur- / Systemdiagramme",
            "more-product-demo": "Mehr Produktlaunch- / Demo-Gefühl",
            "denser-internal": "Dichterer Stil für interne Reviews",
            "custom": "Benutzerdefiniert",
        },
        "visual_expression": {
            "keep": "Aktuell beibehalten",
            "more-premium": "Hochwertiger / stärker gestaltet",
            "more-restrained": "Zurückhaltender / weniger dekorativ",
            "less-cards": "Weniger Karten, offenere Komposition",
            "stronger-evidence": "Stärker evidenzgeführtes Storytelling",
            "custom": "Benutzerdefiniert",
        },
        "image_strategy": {
            "keep": "Aktuell beibehalten",
            "fewer-decorative": "Weniger dekorative Bilder verwenden",
            "more-ai-concept": "KI-Konzeptbilder ergänzen",
            "official-only": "Nur echte Screenshots / offizielle Materialien verwenden",
            "stronger-cover-section": "Titel- oder Kapitelvisuals stärken",
            "custom": "Benutzerdefiniert",
        },
        "logo_strategy": {
            "keep": "Aktuell beibehalten",
            "cover-final-only": "Logos nur auf Titel- und Schlussfolie behalten",
            "footer-brand": "Footer-Branding stärken",
            "partner-logo-page": "Partner- / Schul- / Unternehmenslogo-Seite hinzufügen",
            "remove-all": "Alle Logos entfernen",
            "custom": "Benutzerdefiniert",
        },
    },
    "fr": {
        "output_format": {
            "html-revealjs": "HTML (Reveal.js)",
            "pptx": "PPTX (PowerPoint)",
            "both": "HTML + PPTX (Les deux)",
        },
    },
    "it": {
        "output_format": {
            "html-revealjs": "HTML (Reveal.js)",
            "pptx": "PPTX (PowerPoint)",
            "both": "HTML + PPTX (Entrambi)",
        },
    },
    "es": {
        "output_format": {
            "html-revealjs": "HTML (Reveal.js)",
            "pptx": "PPTX (PowerPoint)",
            "both": "HTML + PPTX (Ambos)",
        },
    },
}

for language, title_items in ADDITIONAL_QUESTION_TITLE_L10N.items():
    QUESTION_TITLE_L10N.setdefault(language, {}).update(title_items)

for language, label_groups in ADDITIONAL_CHOICE_LABEL_L10N.items():
    CHOICE_LABEL_L10N.setdefault(language, {}).update(label_groups)

GENERIC_VISUAL_FIELD_L10N: dict[str, dict[str, str]] = {
    "en": {
        "summary": "Selected visual direction based on the topic, audience, and deck type.",
        "best_for": "A fitting direction for this topic, audience, and deck type.",
        "background": "Use a coherent background system that matches the selected palette.",
        "layout": "Use slide layouts that fit the proof objects and presentation rhythm.",
        "chart": "Use clear evidence-led charts with direct labels.",
        "image_strategy": "Use verified source material; use AI imagery only when authorized.",
        "inspiration": "Use the selected direction as inspiration, not as a rigid template.",
        "risk": "Keep source labels, evidence strength, and layout QA visible.",
    },
    "de": {
        "summary": "Ausgewählte visuelle Richtung auf Basis von Thema, Zielgruppe und Deck-Typ.",
        "best_for": "Eine passende Richtung für dieses Thema, diese Zielgruppe und diesen Deck-Typ.",
        "background": "Ein konsistentes Hintergrundsystem verwenden, das zur gewählten Farbpalette passt.",
        "layout": "Folienlayouts an Beweisobjekte und Präsentationsrhythmus anpassen.",
        "chart": "Klare, evidenzorientierte Diagramme mit direkten Beschriftungen verwenden.",
        "image_strategy": "Geprüftes Quellenmaterial verwenden; KI-Bilder nur mit Freigabe einsetzen.",
        "inspiration": "Die gewählte Richtung als Inspiration nutzen, nicht als starre Vorlage.",
        "risk": "Quellenhinweise, Evidenzstärke und Layout-QA sichtbar halten.",
    },
}


@dataclass(frozen=True)
class Choice:
    value: str
    label: str
    description: str


@dataclass(frozen=True)
class Question:
    key: str
    title: str
    prompt: str
    choices: tuple[Choice, ...]
    default: str


@dataclass(frozen=True)
class VisualCandidate:
    key: str
    name: str
    summary: str
    best_for: str
    avoid_for: str
    palette: tuple[str, str, str, str]
    background: str
    typography: str
    layout: str
    chart: str
    image_strategy: str
    inspiration: str
    risk: str
    html_transition: str = "slide"
    html_animation: str = "minimal"
    html_gradient: str = ""


INTAKE_QUESTIONS: tuple[Question, ...] = (
    Question(
        key="deck_type",
        title="PPT 类型",
        prompt="你要做哪类 PPT?",
        default="engineering-platform",
        choices=(
            Choice("project-report", "项目汇报", "说明进展、结果、风险和下一步。"),
            Choice("engineering-platform", "工程 / 技术方案介绍", "解释系统价值、架构和实现路径。"),
            Choice("investor-pitch", "投资人 / 路演 deck", "突出机会、增长、证明和决策请求。"),
            Choice("knowledge-teaching", "学术 / 课程 / 知识讲解", "把复杂材料重组成清晰知识结构。"),
            Choice("sales-product", "客户销售 / 产品介绍", "展示痛点、方案、demo 和价值证明。"),
            Choice("custom", "自定义", "我有自己的类型描述。"),
        ),
    ),
    Question(
        key="output_format",
        title="输出格式",
        prompt="生成哪种格式的演示文稿？",
        default="html-revealjs",
        choices=(
            Choice(
                "html-revealjs",
                "HTML（Reveal.js）",
                "演示场景首选：动画过渡、presenter mode、浏览器即用，支持 ?print-pdf 导出。",
            ),
            Choice(
                "pptx",
                "PPTX（PowerPoint）",
                "需要可编辑交付、或对方要求 PowerPoint 格式时使用。",
            ),
            Choice(
                "both",
                "HTML + PPTX（两者都生成）",
                "HTML 用于演示，PPTX 用于编辑分享。两版视觉风格会有差异：HTML 版使用渐变背景和动画，PPTX 版使用相同调色板的纯色背景。",
            ),
        ),
    ),
    Question(
        key="research_strategy",
        title="资料研究策略",
        prompt="如果资料不完整，先怎么获得研究材料?",
        default="hybrid-deep-research",
        choices=(
            Choice("hybrid-deep-research", "Hybrid：外部 Deep Research + Codex 定点核验", "适合医学、产业、政策、技术趋势等资料密集主题。"),
            Choice("codex-web-deep", "Codex 深度联网研究", "由 Codex 查找、筛选、核验资料；更耗时和 token。"),
            Choice("external-deep-research", "Gemini / Perplexity Deep Research 资料包", "先用外部 Deep Research 做广覆盖，再把资料包交给 Codex 结构化。"),
            Choice("provided-materials", "只用我提供的资料", "不主动联网，缺失信息必须标注。"),
            Choice("custom", "自定义", "我有自己的研究资料策略。"),
        ),
    ),
    Question(
        key="audience",
        title="听众",
        prompt="这份 PPT 主要给谁看?",
        default="technical-leaders",
        choices=(
            Choice("executives", "高层 / 老板 / 决策者", "更重结论、风险和决策请求。"),
            Choice("investors-reviewers", "投资人 / 评委 / 路演对象", "更重可信证明、市场和增长叙事。"),
            Choice("technical-leaders", "技术团队 / 工程评审", "更重架构、实现、指标和 tradeoff。"),
            Choice("customers-sales", "客户 / 销售对象", "更重痛点、方案、案例和转化。"),
            Choice("teachers-researchers", "学生 / 老师 / 研究者", "更重解释、引用和知识结构。"),
            Choice("custom", "自定义", "我有自己的听众描述。"),
        ),
    ),
    Question(
        key="goal",
        title="目标",
        prompt="这份 PPT 的主要目标是什么?",
        default="explain-value",
        choices=(
            Choice("understand-topic", "让对方快速理解一个主题", "强调清晰解释和结构化。"),
            Choice("decision", "说服对方做决定", "强调证据、取舍和行动请求。"),
            Choice("progress-risk", "汇报进展、成果和风险", "强调状态、指标、风险和下一步。"),
            Choice("teaching", "教学讲解 / 知识传达", "强调概念、例子和学习路径。"),
            Choice("explain-value", "展示产品 / 项目价值", "强调问题、方案、价值和证明。"),
            Choice("custom", "自定义", "我有自己的目标描述。"),
        ),
    ),
    Question(
        key="source_boundary",
        title="资料边界",
        prompt="资料应该怎么使用?",
        default="provided-only",
        choices=(
            Choice("provided-only", "严格只用我提供的材料", "缺失信息必须标注，不联网补全。"),
            Choice("web-with-sources", "可以联网补充，但必须标注来源", "适合需要最新资料或外部事实。"),
            Choice("existing-doc", "以已有 PPT / 文档为内容基础", "继承已有内容结构并改进表达。"),
            Choice("reference-quality", "以参考 deck 作为质量和风格标杆", "参考是质量 bar，不盲目复制。"),
            Choice("custom", "自定义", "我有自己的资料使用规则。"),
        ),
    ),
    Question(
        key="content_language",
        title="内容语言",
        prompt="PPT 正文使用什么语言?",
        default="zh",
        choices=(
            Choice("zh", "中文", "使用中文正文和中文标题。"),
            Choice("en", "English", "Use English slide copy and titles."),
            Choice("de", "Deutsch", "Deutsche Folientexte und Titel verwenden."),
            Choice("fr", "Français", "Utiliser le français pour les titres et le contenu."),
            Choice("it", "Italiano", "Usa l'italiano per titoli e contenuti."),
            Choice("es", "Español", "Usar español en títulos y contenido."),
            Choice("bilingual", "双语", "适合跨语言材料；具体语言可在说明中写明。"),
            Choice("custom", "自定义", "我有自己的语言要求。"),
        ),
    ),
    Question(
        key="output_constraints",
        title="输出限制",
        prompt="页数和演讲时长限制是什么?",
        default="pages-10-12",
        choices=(
            Choice("pages-8-10", "8-10 页", "适合短讲、快速方案或 pitch 初稿。"),
            Choice("pages-10-12", "10-12 页", "默认推荐，适合 10-15 分钟演讲。"),
            Choice("pages-15-20", "15-20 页", "适合详细汇报、内部评审或长材料。"),
            Choice("custom", "自定义时长和页数", "我有自己的页数、时长或结构限制。"),
        ),
    ),
    Question(
        key="logo_policy",
        title="Logo / 品牌素材",
        prompt="Logo 和品牌素材怎么处理?",
        default="provided-only",
        choices=(
            Choice("none", "不使用 logo", "避免品牌资产风险。"),
            Choice("provided-only", "只使用我提供的 logo / 图片", "最安全。"),
            Choice("official-sources", "可以查找官方 logo 和官方素材", "需要记录来源。"),
            Choice("cover-final-only", "只在封面和结束页使用 logo", "弱品牌露出。"),
            Choice("custom", "自定义", "我有自己的品牌素材规则。"),
        ),
    ),
    Question(
        key="image_policy",
        title="AI 生图",
        prompt="是否允许使用 AI 生图?",
        default="ask-before-use",
        choices=(
            Choice("none", "不使用 AI 生成图片", "只用文字、图表、真实素材。"),
            Choice("abstract-only", "只允许生成抽象背景 / 概念图", "不伪造真实对象。"),
            Choice("cover-section", "允许生成封面和章节图", "用于增强视觉表现。"),
            Choice("ask-before-use", "每次生图前先问我", "默认安全策略。"),
            Choice("custom", "自定义", "我有自己的图片策略。"),
        ),
    ),
    Question(
        key="visual_freedom",
        title="第一版视觉方向",
        prompt="第一版视觉方向怎么处理?",
        default="delegate",
        choices=(
            Choice("delegate", "AI 自主决策", "由生成引擎根据主题、听众和视觉候选自主选择最优方案。Codex 环境下由 Presentations 插件主导；Claude Code 环境下由 AI 根据 brief 自主决策。"),
            Choice("restrained", "更正式克制", "适合严肃汇报。"),
            Choice("technical", "更科技 / 工程感", "适合技术方案和架构解释。"),
            Choice("investor", "更投资人路演 / 高对比", "适合 pitch 或评审。"),
            Choice("academic-editorial", "更学术 / 编辑风", "适合知识讲解和研究。"),
            Choice("custom", "自定义", "我有自己的视觉方向。"),
        ),
    ),
    Question(
        key="reference_deck",
        title="参考 deck",
        prompt="是否有参考 deck 或风格样张?",
        default="none",
        choices=(
            Choice("none", "没有参考，按内容生成", "由 Presentations 自主设计。"),
            Choice("quality-only", "有参考，但只作为质量标杆", "beat reference, do not clone。"),
            Choice("visual-style", "有参考，需要接近其视觉风格", "继承风格但不盲目复制。"),
            Choice("existing-ppt", "有已有 PPT，需要在它基础上改", "作为 source/template。"),
            Choice("custom", "自定义", "我有自己的参考使用规则。"),
        ),
    ),
)

REVISION_GROUPS: tuple[Question, ...] = (
    Question(
        key="palette_direction",
        title="配色方向",
        prompt="是否要改变配色方向?",
        default="keep",
        choices=(
            Choice("keep", "保持当前", "不改 v1 配色。"),
            Choice("warm-editorial", "暖色纸感 / 编辑风", "更像高质量长文和研究报告。"),
            Choice("engineering-blue-gray", "工程蓝灰 / 技术风", "更冷静、系统、技术。"),
            Choice("high-contrast-pitch", "高对比 / 路演风", "更强冲击力。"),
            Choice("brand-boost", "强化品牌色", "更明显使用品牌色。"),
            Choice("custom", "自定义", "我有自己的配色方向。"),
        ),
    ),
    Question(
        key="structure_direction",
        title="结构节奏",
        prompt="是否要改变结构节奏?",
        default="keep",
        choices=(
            Choice("keep", "保持当前", "不改整体节奏。"),
            Choice("stronger-story", "更强故事节奏", "章节推进更明显。"),
            Choice("more-analytical", "更数据分析", "更多 evidence-led 页面。"),
            Choice("more-architecture", "更架构 / 系统图", "强化系统图和流程图。"),
            Choice("more-product-demo", "更产品发布 / demo 感", "更像产品展示。"),
            Choice("denser-internal", "更高密度内部汇报", "更适合内部评审。"),
            Choice("custom", "自定义", "我有自己的结构方向。"),
        ),
    ),
    Question(
        key="visual_expression",
        title="视觉表现力",
        prompt="是否要改变视觉表现力?",
        default="keep",
        choices=(
            Choice("keep", "保持当前", "不改表现力。"),
            Choice("more-premium", "更高级 / 更有设计感", "提高整体 polished 感。"),
            Choice("more-restrained", "更克制 / 更少装饰", "减少视觉噪音。"),
            Choice("less-cards", "更少卡片，更开放构图", "避免 generic SaaS card grid。"),
            Choice("stronger-evidence", "更强图表和 evidence-led storytelling", "让证据更直接证明标题。"),
            Choice("custom", "自定义", "我有自己的表现力要求。"),
        ),
    ),
    Question(
        key="image_strategy",
        title="图片策略",
        prompt="图片和视觉素材怎么调整?",
        default="keep",
        choices=(
            Choice("keep", "保持当前", "不改图片策略。"),
            Choice("fewer-decorative", "减少装饰图", "更依赖图表和文字。"),
            Choice("more-ai-concept", "增加 AI 概念图", "只用于已授权场景。"),
            Choice("official-only", "只用真实截图 / 官方素材", "避免 AI 和不明来源素材。"),
            Choice("stronger-cover-section", "增强封面或章节视觉", "提升第一印象。"),
            Choice("custom", "自定义", "我有自己的图片策略。"),
        ),
    ),
    Question(
        key="logo_strategy",
        title="Logo / 品牌露出",
        prompt="Logo 和品牌露出怎么调整?",
        default="keep",
        choices=(
            Choice("keep", "保持当前", "不改品牌露出。"),
            Choice("cover-final-only", "只保留封面和结束页", "降低干扰。"),
            Choice("footer-brand", "增强页脚品牌感", "每页轻品牌提示。"),
            Choice("partner-logo-page", "增加合作方 / 学校 / 企业 logo 页", "适合展示生态或背书。"),
            Choice("remove-all", "移除所有 logo", "避免品牌资产风险。"),
            Choice("custom", "自定义", "我有自己的品牌策略。"),
        ),
    ),
)


def slugify(value: str) -> str:
    cleaned: str = re.sub(r"[^a-zA-Z0-9_-]+", "-", value.strip().lower()).strip("-")
    if cleaned:
        return cleaned[:80]
    return f"presentation-{datetime.now().strftime('%Y%m%d-%H%M%S')}"


def now_id() -> str:
    return datetime.now().strftime("manual-%Y%m%d-%H%M%S")


def workspace_root(base_dir: Path, thread_id: str, task_slug: str) -> Path:
    # `thread_id` is kept for command compatibility. User-facing PPT assets
    # should live together under one project-local folder.
    return base_dir / "PPTX" / task_slug


def status_dir(task_dir: Path) -> Path:
    return task_dir / "status"


def read_json(path: Path, default: JsonDict | None = None) -> JsonDict:
    if not path.exists():
        return {} if default is None else default
    with path.open("r", encoding="utf-8") as handle:
        value: Any = json.load(handle)
    if not isinstance(value, dict):
        raise ValueError(f"Expected object JSON at {path}")
    return value


def write_json(path: Path, data: JsonDict) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    with path.open("w", encoding="utf-8") as handle:
        json.dump(data, handle, ensure_ascii=False, indent=2)
        handle.write("\n")


def write_text(path: Path, text: str) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(text, encoding="utf-8")


def detect_ui_language(text: str) -> str:
    lowered: str = text.lower()
    if re.search(r"[\u4e00-\u9fff]", text):
        return "zh"
    language_markers: dict[str, tuple[str, ...]] = {
        "de": (" der ", " die ", " das ", " und ", " mit ", " für ", "über", " bitte ", " erstellen ", "präsentation"),
        "fr": (" le ", " la ", " les ", " des ", " avec ", " pour ", "présentation", "diapositive", "veuillez"),
        "it": (" il ", " la ", " gli ", " con ", " per ", "presentazione", "diapositiva", "crea"),
        "es": (" el ", " la ", " los ", " con ", " para ", "presentación", "diapositiva", "crear"),
        "en": (" the ", " and ", " with ", " for ", "please ", "create ", "presentation", "slides", "deck"),
    }
    padded: str = f" {lowered} "
    if re.search(r"[äöüß]", lowered):
        return "de"
    if re.search(r"[àâçéèêëîïôûùüÿœ]", lowered):
        return "fr"
    if re.search(r"[áéíóúñ¿¡]", lowered):
        return "es"
    scores: dict[str, int] = {
        lang: sum(1 for marker in markers if marker in padded)
        for lang, markers in language_markers.items()
    }
    best_lang: str = max(scores, key=scores.get)
    return best_lang if scores[best_lang] > 0 else "zh"


def resolve_ui_language(requested: str, conversation_text: str, fallback_text: str) -> str:
    if requested in SUPPORTED_UI_LANGUAGES:
        return requested
    detected: str = detect_ui_language(conversation_text.strip() or fallback_text)
    return detected if detected in SUPPORTED_UI_LANGUAGES else "zh"


def ui_language_from_brief(brief: JsonDict) -> str:
    language: str = str(brief.get("ui_language", "")).strip()
    if language in SUPPORTED_UI_LANGUAGES:
        return language
    fallback_text: str = f"{brief.get('conversation_text', '')} {brief.get('topic', '')}"
    return detect_ui_language(fallback_text)


def ui_language_for_task(task_dir: Path) -> str:
    selected: JsonDict = read_json(task_dir / "intake-selection.json")
    if selected:
        return ui_language_from_brief(selected)
    return ui_language_from_brief(read_json(task_dir / "brief-draft.json"))


def t(ui_language: str, key: str) -> str:
    language_copy: dict[str, str] = UI_COPY.get(ui_language, {})
    if key in language_copy:
        return language_copy[key]
    if ui_language != "zh" and key in UI_COPY["en"]:
        return UI_COPY["en"][key]
    return UI_COPY["zh"].get(key, key)


def generation_strategy_text(output_format: str, task_dir: Path, ui_language: str) -> str:
    task_path: str = str(task_dir)
    messages: dict[str, dict[str, str]] = {
        "zh": {
            "html-revealjs": f"将生成 Reveal.js HTML 演示文稿，保存到 {task_path}/final/<name>.html。用浏览器打开即可演示；附加 ?print-pdf 可导出 PDF。",
            "pptx": f"先生成 v1 PPTX 和 contact sheet，集中保存到 {task_path}，然后打开 style-review.html 供选择是否重绘。",
            "both": f"将分别生成 HTML（Reveal.js）和 PPTX 两个版本，均保存到 {task_path}/final/。两版视觉风格会有差异：HTML 使用渐变背景和动画，PPTX 使用纯色背景。",
        },
        "en": {
            "html-revealjs": f"Generate a Reveal.js HTML presentation and save it to {task_path}/final/<name>.html. Open it in a browser to present; append ?print-pdf to export PDF.",
            "pptx": f"Generate the v1 PPTX and contact sheet first, save them under {task_path}, then open style-review.html so you can decide whether to redraw the deck.",
            "both": f"Generate both HTML (Reveal.js) and PPTX versions under {task_path}/final/. The visual systems intentionally differ: HTML uses gradients and animation, while PPTX uses solid-color equivalents.",
        },
        "de": {
            "html-revealjs": f"Es wird eine Reveal.js-HTML-Präsentation erzeugt und unter {task_path}/final/<name>.html gespeichert. Im Browser öffnen; mit ?print-pdf als PDF exportieren.",
            "pptx": f"Zuerst werden v1-PPTX und Contact Sheet erzeugt und unter {task_path} gespeichert. Danach wird style-review.html geöffnet, damit Sie entscheiden können, ob das Deck visuell überarbeitet werden soll.",
            "both": f"Es werden HTML (Reveal.js) und PPTX unter {task_path}/final/ erzeugt. Die visuellen Systeme unterscheiden sich bewusst: HTML nutzt Verläufe und Animationen, PPTX nutzt entsprechende Vollfarben.",
        },
        "fr": {
            "html-revealjs": f"Générer une présentation HTML Reveal.js dans {task_path}/final/<name>.html. Ouvrez-la dans un navigateur; ajoutez ?print-pdf pour exporter en PDF.",
            "pptx": f"Générer d'abord le PPTX v1 et la planche de contact dans {task_path}, puis ouvrir style-review.html pour décider d'une éventuelle refonte visuelle.",
            "both": f"Générer les versions HTML (Reveal.js) et PPTX dans {task_path}/final/. Les styles diffèrent volontairement: HTML utilise des dégradés et animations, PPTX utilise des aplats équivalents.",
        },
        "it": {
            "html-revealjs": f"Genera una presentazione HTML Reveal.js in {task_path}/final/<name>.html. Aprila nel browser; aggiungi ?print-pdf per esportare in PDF.",
            "pptx": f"Genera prima il PPTX v1 e il contact sheet in {task_path}, poi apri style-review.html per decidere se ridisegnare il deck.",
            "both": f"Genera sia HTML (Reveal.js) sia PPTX in {task_path}/final/. Gli stili differiscono intenzionalmente: HTML usa gradienti e animazioni, PPTX usa colori pieni equivalenti.",
        },
        "es": {
            "html-revealjs": f"Genera una presentación HTML Reveal.js en {task_path}/final/<name>.html. Ábrela en el navegador; añade ?print-pdf para exportar a PDF.",
            "pptx": f"Primero genera el PPTX v1 y la hoja de contacto en {task_path}, y luego abre style-review.html para decidir si redibujar el deck.",
            "both": f"Genera versiones HTML (Reveal.js) y PPTX en {task_path}/final/. Los estilos difieren a propósito: HTML usa degradados y animación, PPTX usa colores sólidos equivalentes.",
        },
    }
    language_messages: dict[str, str] = messages.get(ui_language, messages["en"])
    return language_messages.get(output_format, language_messages["pptx"])


def localized_question_title(question: Question, ui_language: str) -> str:
    if ui_language == "zh":
        return question.title
    return (
        QUESTION_TITLE_L10N.get(ui_language, {}).get(question.key)
        or QUESTION_TITLE_L10N.get("en", {}).get(question.key)
        or humanize_key(question.key)
    )


def localized_question_prompt(question: Question, ui_language: str) -> str:
    if ui_language == "zh":
        return question.prompt
    return (
        QUESTION_PROMPT_L10N.get(ui_language, {}).get(question.key)
        or QUESTION_PROMPT_L10N.get("en", {}).get(question.key)
        or ""
    )


def humanize_key(value: str) -> str:
    return value.replace("_", " ").replace("-", " ").strip().title()


def localized_choice_label_value(question: Question, value: str, fallback: str, ui_language: str) -> str:
    if ui_language == "zh":
        return fallback
    return (
        CHOICE_LABEL_L10N.get(ui_language, {}).get(question.key, {}).get(value)
        or CHOICE_LABEL_L10N.get("en", {}).get(question.key, {}).get(value)
        or humanize_key(value)
    )


def localized_choice_label(question: Question, item: JsonDict, ui_language: str) -> str:
    value: str = str(item.get("value", ""))
    if "custom" in item:
        return str(item.get("custom", item.get("label", "")))
    return localized_choice_label_value(question, value, str(item.get("label", "")), ui_language)


def localized_choice_description(question: Question, choice: Choice, ui_language: str) -> str:
    if ui_language == "zh":
        return choice.description
    return ""


def localized_source(source_name: str, ui_language: str) -> str:
    if source_name == "default":
        return t(ui_language, "default")
    if source_name == "user-selected":
        return t(ui_language, "user_selected")
    if not source_name:
        return t(ui_language, "unknown")
    return source_name


def localized_risk(risk: str, ui_language: str) -> str:
    if ui_language == "zh":
        return risk
    risk_map: dict[str, dict[str, str]] = {
        "未发现明确 logo 文件；如需使用 logo，必须由用户提供或使用官方来源。": {
            "en": "No explicit logo file was found; if logos are needed, they must be provided by the user or sourced officially.",
            "de": "Keine eindeutige Logo-Datei gefunden; falls Logos benötigt werden, müssen sie vom Benutzer bereitgestellt oder aus offiziellen Quellen bezogen werden.",
        },
        "未发现明确量化数据文件；第一版可能需要用定性证明或标注缺失指标。": {
            "en": "No explicit quantitative data file was found; v1 may need qualitative evidence or clear missing-metric labels.",
            "de": "Keine eindeutige quantitative Datendatei gefunden; v1 benötigt eventuell qualitative Evidenz oder klare Hinweise auf fehlende Kennzahlen.",
        },
        "未提供具体资料路径；需要在生成前补充 source material。": {
            "en": "No concrete source path was provided; source material should be added before generation.",
            "de": "Es wurde kein konkreter Quellenpfad angegeben; Quellenmaterial sollte vor der Generierung ergänzt werden.",
        },
    }
    return risk_map.get(risk, {}).get(ui_language, risk_map.get(risk, {}).get("en", risk))


def localized_visual_field(candidate: JsonDict, field: str, ui_language: str) -> str:
    value: str = str(candidate.get(field, ""))
    if ui_language == "zh":
        return value
    return GENERIC_VISUAL_FIELD_L10N.get(ui_language, GENERIC_VISUAL_FIELD_L10N.get("en", {})).get(field, value)


def natural_sort_key(path: Path) -> list[tuple[int, int | str]]:
    parts: list[str] = re.split(r"(\d+)", path.name.lower())
    return [(0, int(part)) if part.isdigit() else (1, part) for part in parts]


def touch_status(task_dir: Path, status_name: str) -> Path:
    filename: str | None = STATUS_FILES.get(status_name)
    if filename is None:
        raise ValueError(f"Unknown status: {status_name}")
    path: Path = status_dir(task_dir) / filename
    write_text(path, datetime.now().isoformat(timespec="seconds") + "\n")
    return path


def confirm_token_path(task_dir: Path) -> Path:
    return status_dir(task_dir) / "confirm.token"


def ensure_confirm_token(task_dir: Path) -> str:
    path: Path = confirm_token_path(task_dir)
    if path.exists():
        token: str = path.read_text(encoding="utf-8").strip()
        if token:
            return token
    token = secrets.token_urlsafe(24)
    write_text(path, token + "\n")
    return token


def valid_confirm_token(task_dir: Path, token: str) -> bool:
    path: Path = confirm_token_path(task_dir)
    if not path.exists() or not token:
        return False
    return secrets.compare_digest(path.read_text(encoding="utf-8").strip(), token)


def confirmation_receipt(token: str, confirmed_at: str) -> JsonDict:
    return {
        "method": "browser-form",
        "confirmed_by": "user-click",
        "token_verified": True,
        "token_sha256": hashlib.sha256(token.encode("utf-8")).hexdigest(),
        "confirmed_at": confirmed_at,
    }


def validate_generation_guard(task_dir: Path) -> list[str]:
    errors: list[str] = []
    brief: JsonDict = read_json(task_dir / "brief-confirmed.json")
    if not brief:
        errors.append(f"Missing confirmed brief: {task_dir / 'brief-confirmed.json'}")
        return errors
    if brief.get("confirmed") is not True:
        errors.append("Confirmed brief exists but confirmed is not true.")
    if not (status_dir(task_dir) / STATUS_FILES["confirmed"]).exists():
        errors.append(f"Missing confirmation status: {status_dir(task_dir) / STATUS_FILES['confirmed']}")
    receipt: Any = brief.get("confirmation_gate")
    if not isinstance(receipt, dict):
        errors.append("Missing confirmation_gate receipt; open the confirmation page and submit the form.")
    else:
        if receipt.get("method") != "browser-form":
            errors.append("confirmation_gate.method is not browser-form.")
        if receipt.get("confirmed_by") != "user-click":
            errors.append("confirmation_gate.confirmed_by is not user-click.")
        if receipt.get("token_verified") is not True:
            errors.append("confirmation_gate.token_verified is not true.")
    return errors


def selected_choice(question: Question, value: str) -> Choice:
    for choice in question.choices:
        if choice.value == value:
            return choice
    return next(choice for choice in question.choices if choice.value == question.default)


def default_intake_value(question: Question, sources: list[str], ui_language: str = "zh") -> str:
    if question.key == "content_language" and ui_language in {"zh", "en", "de", "fr", "it", "es"}:
        return ui_language
    if question.key == "source_boundary":
        return "provided-only" if sources else "web-with-sources"
    if question.key == "research_strategy":
        if not sources:
            return "hybrid-deep-research"
        source_text: str = " ".join(sources).lower()
        if any(token in source_text for token in ("gemini", "perplexity", "deep-research", "deep research")):
            return "external-deep-research"
        return "provided-materials"
    return question.default


def selection_value(selections: JsonDict, key: str, default: str = "") -> str:
    item: Any = selections.get(key, {})
    if isinstance(item, dict):
        return str(item.get("value", default))
    return default


def output_format_from_selections(selections: JsonDict, default: str = "pptx") -> str:
    output_format: str = selection_value(selections, "output_format", default)
    if output_format in {"html-revealjs", "pptx", "both"}:
        return output_format
    return default


def visual_candidate_to_json(candidate: VisualCandidate) -> JsonDict:
    return {
        "key": candidate.key,
        "name": candidate.name,
        "summary": candidate.summary,
        "best_for": candidate.best_for,
        "avoid_for": candidate.avoid_for,
        "palette": list(candidate.palette),
        "background": candidate.background,
        "typography": candidate.typography,
        "layout": candidate.layout,
        "chart": candidate.chart,
        "image_strategy": candidate.image_strategy,
        "inspiration": candidate.inspiration,
        "risk": candidate.risk,
        "html_transition": candidate.html_transition,
        "html_animation": candidate.html_animation,
        "html_gradient": candidate.html_gradient,
    }


def classify_visual_context(topic: str, selections: JsonDict) -> str:
    deck_type: str = selection_value(selections, "deck_type")
    audience: str = selection_value(selections, "audience")
    text: str = f"{topic} {deck_type} {audience}".lower()
    if any(token in text for token in ("体育", "足球", "篮球", "运动", "赛事", "sports", "football", "basketball", "club")):
        return "sports"
    if any(token in text for token in ("战略", "策略", "规划", "运营", "组织", "增长", "strategy", "planning", "operations")):
        return "strategy"
    if any(token in text for token in ("医学", "药", "临床", "疾病", "研究", "alzheimer", "clinical", "biotech", "medical", "research")):
        return "research"
    if deck_type in {"engineering-platform"} or any(token in text for token in ("工程", "架构", "系统", "平台", "infra", "architecture", "developer")):
        return "engineering"
    if deck_type in {"investor-pitch", "sales-product"}:
        return "market"
    return "general"


def html_profile_for_candidate(context: str, candidate: VisualCandidate) -> tuple[str, str, str]:
    text: str = f"{context} {candidate.key} {candidate.name} {candidate.summary}".lower()
    if any(token in text for token in ("pitch", "investor", "studio-pitch", "路演", "投资")):
        return "zoom", "rich", "linear-gradient(135deg, #1a1a2e, #16213e, #0f3460)"
    if any(token in text for token in ("product", "launch", "brand", "studio-visual", "creative", "产品", "创意")):
        return "convex", "rich", "linear-gradient(135deg, #667eea, #764ba2)"
    if "aurora" in text:
        return "zoom", "rich", "linear-gradient(135deg, #e0f2fe, #dbeafe, #f0fdf4)"
    if context == "research" or any(token in text for token in ("academic", "clinical", "medical", "atlas", "research", "学术", "知识")):
        return "fade", "minimal", "linear-gradient(135deg, #f5f7fa, #c3cfe2)"
    if any(token in text for token in ("boardroom", "consulting", "roadmap", "restrained", "正式", "克制")):
        return "fade", "minimal", ""
    if context == "engineering" or any(token in text for token in ("engineering", "terminal", "system", "signal", "architecture", "科技", "工程")):
        return "slide", "moderate", "linear-gradient(135deg, #0f0c29, #302b63, #24243e)"
    return "fade", "minimal", ""


def with_html_profile(context: str, candidate: VisualCandidate) -> VisualCandidate:
    html_transition, html_animation, html_gradient = html_profile_for_candidate(context, candidate)
    return replace(
        candidate,
        html_transition=html_transition,
        html_animation=html_animation,
        html_gradient=html_gradient,
    )


def build_visual_candidates(topic: str, selections: JsonDict) -> tuple[VisualCandidate, ...]:
    context: str = classify_visual_context(topic, selections)
    libraries: dict[str, tuple[VisualCandidate, ...]] = {
        "sports": (
            VisualCandidate(
                "broadcast-analytics",
                "Broadcast Analytics",
                "像体育转播数据大屏，节奏强、数字醒目、对比高。",
                "赛事复盘、球队经营、体育商业分析。",
                "严肃学术汇报或低调内部材料。",
                ("#07111f", "#f7f8fb", "#24d18b", "#ffcc33"),
                "深色场地感背景，少量动线和数据网格。",
                "粗标题 + 紧凑数字标签。",
                "比分板式大数字、左右对抗、关键事件时间线。",
                "排名、雷达、趋势线，直接标注而少用图例。",
                "真实赛事/训练素材优先；没有素材时使用抽象运动轨迹。",
                "sports broadcast UI + html deck theme catalog + evidence dashboard",
                "容易过于娱乐化，商业策略页需要降噪。",
            ),
            VisualCandidate(
                "performance-lab",
                "Performance Lab",
                "运动科学实验室风格，干净、精准、强调指标和方法。",
                "训练表现、运动医学、青训体系、数据分析。",
                "需要强情绪感染力的路演封面。",
                ("#f4f7f8", "#17202a", "#006b7a", "#f05a28"),
                "浅底实验室网格，局部使用测量线和坐标标尺。",
                "清晰无衬线，数字使用 tabular 风格。",
                "指标卡 + 流程图 + 对比切片。",
                "小 multiples、区间带、前后对照。",
                "真实运动员素材需来源明确；可用抽象人体运动线框。",
                "performance dashboard + scientific poster + ui-ux-pro-max chart grammar",
                "可能偏理性，封面需要额外增强视觉冲击。",
            ),
            VisualCandidate(
                "club-strategy-room",
                "Club Strategy Room",
                "俱乐部董事会/战术室风格，适合把体育主题讲成战略问题。",
                "球队战略、商业化、组织治理、赞助规划。",
                "纯技术训练分析。",
                ("#102033", "#f0eadc", "#b99855", "#e85d04"),
                "暖纸 + 深色战术线，背景像会议室白板和策略地图。",
                "稳重标题字体，正文较克制。",
                "战略地图、路径图、stakeholder matrix。",
                "瀑布图、路线图、二维象限。",
                "尽量使用真实 logo/球衣/场馆素材；无授权则不用。",
                "strategy room + editorial deck + locked visual contract",
                "如果资料很多，需要避免把每页做成咨询报告表格。",
            ),
        ),
        "strategy": (
            VisualCandidate(
                "executive-boardroom",
                "Executive Boardroom",
                "高层决策汇报风格，克制但有权威感。",
                "战略规划、预算、组织调整、董事会汇报。",
                "创意发布会或教学课程。",
                ("#0f172a", "#f8fafc", "#3b82f6", "#c8a24a"),
                "浅底为主，关键页使用深色整页强调。",
                "强层级标题，正文短句化。",
                "决策摘要、风险矩阵、路线图、经营仪表盘。",
                "waterfall、scenario matrix、milestone roadmap。",
                "少装饰，优先使用真实业务图表和图标。",
                "consulting deck + executive dashboard + design-lock discipline",
                "如果全程深色，会显得压抑；应控制深色页比例。",
            ),
            VisualCandidate(
                "consulting-roadmap",
                "Consulting Roadmap",
                "咨询公司路线图风格，结构强、阶段清楚。",
                "三年规划、转型项目、市场进入策略。",
                "需要强品牌个性的产品发布。",
                ("#ffffff", "#111827", "#2563eb", "#f97316"),
                "纯白背景，细线分区，少量高饱和强调。",
                "标题结论句化，正文像 executive memo。",
                "阶段推进、2x2、能力地图、优先级排序。",
                "矩阵、堆叠条、战略地图。",
                "不依赖图片，强调结构图和业务对象。",
                "strategy scaffold + html-ppt-skill template catalog + ui-ux-pro-max UX rules",
                "容易变成普通咨询模板，需要用主题专属 proof object 打破通用感。",
            ),
            VisualCandidate(
                "operating-model-dashboard",
                "Operating Model Dashboard",
                "运营模型仪表盘风格，适合把规划落到指标和机制。",
                "年度经营计划、组织运行、OKR、流程治理。",
                "纯品牌故事或情绪型演讲。",
                ("#eef2f6", "#1f2937", "#0f766e", "#7c3aed"),
                "浅灰工作台背景，模块边界清楚，强调可扫描。",
                "中等字号、高信息密度但留白稳定。",
                "机制图、RACI、节奏表、指标看板。",
                "KPI cards、heatmap、pipeline funnel。",
                "真实流程图/组织图优先，不用装饰性背景图。",
                "operational SaaS UI + dashboard grammar + quality gates",
                "过密时容易不适合演讲，需要 split slide。",
            ),
        ),
        "research": (
            VisualCandidate(
                "clinical-atlas",
                "Clinical Atlas",
                "研究图谱风格，适合复杂证据链和研发路径。",
                "医学、生命科学、政策研究、技术综述。",
                "销售型产品发布和强情绪路演。",
                ("#f6f8fb", "#132238", "#0f6f8f", "#d97706"),
                "浅底研究图谱，局部使用细网格、节点和证据线。",
                "学术无衬线，标题清楚，来源小字统一。",
                "证据链、pipeline、机制图、分层框架。",
                "pipeline、forest-like comparison、evidence table。",
                "真实图表和示意图优先；AI 图只作抽象章节背景。",
                "academic design-lock + evidence dashboard + source-first QA",
                "如果每页都很理性，封面和章节页需要更强视觉记忆点。",
            ),
            VisualCandidate(
                "biotech-pipeline",
                "Biotech Pipeline",
                "生物科技投资人 deck 风格，重点突出机会、管线和决策。",
                "药物研发、创新技术、融资路演、BD 汇报。",
                "纯课程教学或极保守学术答辩。",
                ("#08111f", "#f8fafc", "#4cc9f0", "#ffb703"),
                "深浅交替，封面和章节页有高对比抽象分子/路径背景。",
                "大标题 + 突出 milestone 和数据。",
                "研发管线、市场窗口、风险消减、里程碑。",
                "pipeline chart、TAM blocks、risk ladder。",
                "可生成抽象科学背景，但不能伪造实验图片或监管标志。",
                "startup pitch + clinical pipeline layout + visual preview selection",
                "容易过度路演化，需保留来源和监管谨慎语气。",
            ),
            VisualCandidate(
                "medical-editorial",
                "Medical Editorial",
                "高质量医学长文/期刊特稿风格，适合解释复杂主题。",
                "研究综述、课程讲解、公众科普、专家分享。",
                "需要强销售转化的商业 deck。",
                ("#f2efe8", "#1f2933", "#6b7c93", "#b35c1e"),
                "暖纸底、细分隔线、少量图像纹理。",
                "出版物式标题，正文段落可读性优先。",
                "定义、分层、对比、时间线、注释式图表。",
                "annotated timeline、small multiples、callout table。",
                "少量抽象背景；内容页保持干净。",
                "editorial design-lock + beautiful template contract + typography rules",
                "可能不够有冲击力，商业听众需要增加关键数字页。",
            ),
            VisualCandidate(
                "aurora-light",
                "Aurora Light",
                "浅色底+渐变光晕+动感数据展示。兼顾可读性与科技感动画，深色字配发光强调色。",
                "学术讲演、公众演讲、科技科普、混合光源投影场景。",
                "需要强戏剧感的深色大屏路演或暗场展示。",
                ("#f0f9ff", "#0d2137", "#0d7ea2", "#d97706"),
                "浅蓝白底，封面使用淡蓝/薄荷渐变；内容页用渐变色标题和发光数字强调。",
                "深色主文字 + 渐变色大标题（CSS gradient-text）+ 发光数字。",
                "发光数据卡、渐变进度条、动态机制图、3D标题效果。",
                "渐变色填充进度条、光晕数字卡、轻玻璃感卡片。",
                "浅底+淡色渐变，封面使用蓝/薄荷渐变光晕；内容页保持高对比深色文字。",
                "aurora design + glassmorphism light + gradient-text + animated accents",
                "投影机亮度不足时渐变光晕可能变浅，建议调高亮度；避免纯白内容页缺乏层次。",
            ),
        ),
        "engineering": (
            VisualCandidate(
                "signal-system",
                "Signal System",
                "系统评审风格，清晰、工程化、数据前置。",
                "平台架构、工程汇报、技术路线、SRE/数据系统。",
                "品牌发布会或教学科普。",
                ("#f8fafc", "#111827", "#0061ff", "#16a34a"),
                "浅底工程网格，少量蓝色信号线。",
                "紧凑标题，代码/指标用等宽或 tabular 数字。",
                "架构图、数据流、模块责任、指标看板。",
                "latency charts、dependency map、before/after bars。",
                "真实截图优先，截图可美化但不重画事实。",
                "Signal design language + architecture QA + screenshot treatment rules",
                "如果不加故事节奏，容易像技术文档。",
            ),
            VisualCandidate(
                "terminal-architecture",
                "Terminal Architecture",
                "高对比开发者工具风格，适合代码和基础设施。",
                "infra、DevOps、AI agent、开发者平台。",
                "非技术高层或客户销售。",
                ("#09090b", "#fafafa", "#7dd3fc", "#a3e635"),
                "深色终端背景，细线拓扑和命令行提示。",
                "大标题清晰，代码块严格控制长度。",
                "系统拓扑、调用链、状态机、故障演练。",
                "sequence diagram、service map、error budget chart。",
                "避免假 UI；真实终端/产品截图需来源明确。",
                "Terminal design language + locked layout IDs + render QA",
                "不适合长正文，文字必须拆页。",
            ),
            VisualCandidate(
                "product-operations",
                "Product Operations",
                "产品运营工作台风格，适合把技术价值讲给业务方。",
                "产品方案、客户场景、运营效率、自动化平台。",
                "底层架构深挖。",
                ("#f4f6f8", "#1f2937", "#0ea5e9", "#f59e0b"),
                "浅灰 SaaS 工作台，但减少普通卡片堆叠。",
                "清楚的组件标题和短解释。",
                "用户旅程、功能流程、价值闭环、仪表盘。",
                "funnel、cohort trend、task flow。",
                "真实产品截图优先，必要时使用线框场景图。",
                "SaaS UI + product launch scaffold + no generic card grid rule",
                "若卡片过多会显得普通，需要开放构图页做节奏变化。",
            ),
        ),
        "market": (
            VisualCandidate(
                "studio-pitch",
                "Studio Pitch",
                "更像高质量创业路演，节奏快、对比强、视觉记忆点明显。",
                "融资、比赛、产品发布、BD。",
                "监管、医学、法务等需要克制语气的汇报。",
                ("#101828", "#ffffff", "#5b7cfa", "#ff6b35"),
                "深浅切换，大封面和章节图形成记忆点。",
                "短标题、大数字、强对比。",
                "问题-方案-证明-增长-请求。",
                "traction chart、market map、before/after。",
                "可以用抽象概念图，但不能伪造客户 logo 或产品截图。",
                "Studio design language + pitch deck scaffold + visual-led HTML previews",
                "若证据不足，强表现力会放大可信度风险。",
            ),
            VisualCandidate(
                "product-launch",
                "Product Launch",
                "产品发布会风格，突出用户场景和体验。",
                "新品发布、功能介绍、客户演示。",
                "高密度内部复盘。",
                ("#ffffff", "#0f172a", "#7c3aed", "#06b6d4"),
                "干净背景 + 大图区域 + 少量品牌渐变。",
                "标题直接，辅助文案短。",
                "场景、功能、流程、价值证明。",
                "feature matrix、workflow、adoption trend。",
                "真实产品截图优先；无截图时用抽象线框，不伪造 UI。",
                "product launch template + screenshot slot contract + ui-ux-pro-max style match",
                "如果没有真实素材，视觉会偏概念化。",
            ),
            VisualCandidate(
                "brand-narrative",
                "Brand Narrative",
                "品牌叙事风格，强调信任、愿景和故事线。",
                "品牌升级、客户提案、公益/教育项目。",
                "纯技术评审或数据审计。",
                ("#f7f2e8", "#1c1917", "#2563eb", "#c2410c"),
                "暖底、摄影/插画槽位、章节像短篇故事。",
                "叙事标题，正文强调节奏。",
                "用户故事、价值阶梯、生态地图。",
                "impact metrics、quote callout、timeline。",
                "真实人物/客户素材必须授权；否则使用抽象场景图。",
                "editorial template + brand system + visual contract depth",
                "如果目标是做决策，故事页必须配决策页。",
            ),
        ),
    }
    default_candidates: tuple[VisualCandidate, ...] = (
        VisualCandidate(
            "atlas-explanatory",
            "Atlas Explanatory",
            "结构化知识图谱风格，适合把复杂主题讲清楚。",
            "研究、课程、政策、行业分析。",
            "强销售转化或纯代码评审。",
            ("#f6f3ee", "#1f2937", "#2563eb", "#b45309"),
            "浅底、细网格、少量章节色块。",
            "稳健标题，正文强调可读性。",
            "框架图、时间线、对比表、证据链。",
            "timeline、matrix、annotated chart。",
            "抽象背景只用于封面和章节页。",
            "Atlas design language + design-lock registry + source-first planner",
            "可能偏解释型，需要在封面增强主题记忆点。",
        ),
        VisualCandidate(
            "signal-dashboard",
            "Signal Dashboard",
            "数据和机制先行，适合把主题变成可行动判断。",
            "策略、运营、工程、管理汇报。",
            "情绪型演讲或品牌故事。",
            ("#f8fafc", "#0f172a", "#0f766e", "#f59e0b"),
            "浅灰工作台背景，图表和流程为主。",
            "数字清楚，标签短。",
            "dashboard、flow、roadmap、decision matrix。",
            "KPI、trend、heatmap、funnel。",
            "少图片，多使用原生图表和形状。",
            "Signal design language + dashboard UI + quality gates",
            "容易过密，必须控制每页 proof object 数量。",
        ),
        VisualCandidate(
            "studio-visual",
            "Studio Visual",
            "更强视觉表达，适合需要第一眼吸引力的 deck。",
            "路演、发布、竞赛、客户提案。",
            "保守审计或纯学术材料。",
            ("#111827", "#ffffff", "#6366f1", "#f97316"),
            "强封面、章节视觉和深浅对比。",
            "大标题、短句、视觉锚点。",
            "hero statement、before/after、proof snapshot。",
            "大数字、对比图、简化趋势。",
            "允许抽象概念图，禁止伪造真实资产。",
            "Studio design language + 3-preview workflow + visual-led templates",
            "表现力强时更需要事实来源压住可信度。",
        ),
    )
    selected_candidates: tuple[VisualCandidate, ...] = libraries.get(context, default_candidates)
    return tuple(with_html_profile(context, candidate) for candidate in selected_candidates)


def html_page(title: str, body: str, ui_language: str = "zh") -> str:
    html_language: str = HTML_LANG.get(ui_language, HTML_LANG["zh"])
    return f"""<!doctype html>
<html lang="{html.escape(html_language)}">
<head>
  <meta charset="utf-8">
  <meta name="viewport" content="width=device-width, initial-scale=1">
  <title>{html.escape(title)}</title>
  <style>
    :root {{
      --bg: #f7f4ee;
      --panel: #ffffff;
      --ink: #161616;
      --muted: #626262;
      --line: #d8d2c7;
      --accent: #274c77;
      --accent-2: #c75000;
    }}
    * {{ box-sizing: border-box; }}
    body {{
      margin: 0;
      background: var(--bg);
      color: var(--ink);
      font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", "Noto Sans SC", sans-serif;
      line-height: 1.5;
    }}
    main {{ max-width: 1120px; margin: 0 auto; padding: 40px 24px 72px; }}
    h1 {{ font-size: 32px; margin: 0 0 8px; letter-spacing: 0; }}
    h2 {{ font-size: 22px; margin: 28px 0 12px; }}
    h3 {{ font-size: 17px; margin: 0 0 8px; }}
    p {{ color: var(--muted); margin: 0 0 16px; }}
    .topline {{ text-transform: uppercase; letter-spacing: .08em; color: var(--accent-2); font-size: 12px; font-weight: 700; }}
    .grid {{ display: grid; grid-template-columns: repeat(auto-fit, minmax(240px, 1fr)); gap: 12px; }}
    .card, .section {{
      background: var(--panel);
      border: 1px solid var(--line);
      border-radius: 8px;
      padding: 18px;
      box-shadow: 0 1px 0 rgba(0,0,0,.03);
    }}
    .section {{ margin-top: 18px; }}
    label.option {{
      display: block;
      border: 1px solid var(--line);
      border-radius: 8px;
      padding: 14px;
      cursor: pointer;
      background: #fff;
      min-height: 112px;
    }}
    label.option:has(input:checked) {{
      border-color: var(--accent);
      box-shadow: 0 0 0 2px rgba(39, 76, 119, .16);
    }}
    input[type="radio"] {{ margin-right: 8px; }}
    input[type="text"], textarea {{
      width: 100%;
      margin-top: 8px;
      padding: 10px 12px;
      border: 1px solid var(--line);
      border-radius: 6px;
      font: inherit;
    }}
    textarea {{ min-height: 88px; resize: vertical; }}
    .actions {{ display: flex; gap: 12px; flex-wrap: wrap; margin-top: 28px; }}
    button, .button {{
      appearance: none;
      border: 0;
      background: var(--accent);
      color: #fff;
      padding: 12px 18px;
      border-radius: 6px;
      font-weight: 700;
      cursor: pointer;
      text-decoration: none;
      display: inline-block;
    }}
    button.secondary, .button.secondary {{ background: #545454; }}
    button.warning {{ background: var(--accent-2); }}
    .meta {{ color: var(--muted); font-size: 13px; }}
    .source-tag {{ display: inline-block; padding: 2px 7px; border-radius: 999px; background: #ebe3d5; font-size: 12px; color: #554; }}
    .risk {{ border-left: 4px solid var(--accent-2); padding-left: 12px; }}
    .candidate {{ min-height: 100%; }}
    .candidate input {{ margin-right: 8px; }}
    .swatches-preview {{ display: flex; gap: 6px; margin: 12px 0; }}
    .swatch-preview {{ width: 34px; height: 22px; border-radius: 4px; border: 1px solid rgba(0,0,0,.12); }}
    .mini-slides {{ display: grid; grid-template-columns: 1fr 1fr; gap: 8px; margin: 12px 0; }}
    .mini-slide {{
      aspect-ratio: 16 / 9;
      border-radius: 6px;
      padding: 10px;
      display: flex;
      flex-direction: column;
      justify-content: space-between;
      border: 1px solid rgba(0,0,0,.14);
      overflow: hidden;
    }}
    .mini-slide strong {{ font-size: 12px; line-height: 1.15; }}
    .mini-slide span {{ display: block; width: 60%; height: 5px; border-radius: 99px; opacity: .9; }}
    .mini-bars {{ display: grid; gap: 4px; }}
    .mini-bars i {{ display: block; height: 5px; border-radius: 99px; opacity: .86; }}
    .html-field {{ display: flex; gap: 8px; align-items: center; margin: 6px 0; font-size: 13px; }}
    .html-field .label {{ color: var(--muted); min-width: 86px; }}
    .html-field .value {{ color: var(--ink); overflow-wrap: anywhere; }}
    .contact-sheet {{
      width: 100%;
      max-height: 520px;
      object-fit: contain;
      border: 1px solid var(--line);
      background: #eee;
      border-radius: 8px;
    }}
    code {{ background: #eee7da; padding: 2px 5px; border-radius: 4px; }}
    table {{ width: 100%; border-collapse: collapse; background: var(--panel); }}
    th, td {{ border: 1px solid var(--line); padding: 10px; text-align: left; vertical-align: top; }}
    th {{ background: #efe8dc; }}
  </style>
</head>
<body>
<main>
{body}
</main>
</body>
</html>
"""


def question_section(question: Question, current: str = "", ui_language: str = "zh") -> str:
    current_value: str = current or question.default
    cards: list[str] = []
    for choice in question.choices:
        checked: str = " checked" if choice.value == current_value else ""
        custom_input: str = ""
        if choice.value == "custom":
            custom_input = (
                f'<input type="text" name="{question.key}__custom" '
                f'placeholder="{html.escape(t(ui_language, "custom_placeholder"))}">'
            )
        description: str = localized_choice_description(question, choice, ui_language)
        description_html: str = f"<p>{html.escape(description)}</p>" if description else ""
        cards.append(
            f"""<label class="option">
  <input type="radio" name="{question.key}" value="{html.escape(choice.value)}"{checked}>
  <strong>{html.escape(localized_choice_label_value(question, choice.value, choice.label, ui_language))}</strong>
  {description_html}
  {custom_input}
</label>"""
        )
    return f"""<section class="section">
  <h2>{html.escape(localized_question_title(question, ui_language))}</h2>
  <p>{html.escape(localized_question_prompt(question, ui_language))}</p>
  <div class="grid">
    {''.join(cards)}
  </div>
</section>"""


def build_draft_brief(
    task_slug: str,
    topic: str,
    sources: list[str],
    ui_language: str = "auto",
    conversation_text: str = "",
    enhance_mode: bool = False,
) -> JsonDict:
    source_items: list[JsonDict] = build_source_items(sources)
    fallback_text: str = " ".join([topic, task_slug, *sources])
    resolved_ui_language: str = resolve_ui_language(ui_language, conversation_text, fallback_text)
    default_selections: JsonDict = {
        question.key: {
            "value": default_intake_value(question, sources, resolved_ui_language),
            "label": selected_choice(question, default_intake_value(question, sources, resolved_ui_language)).label,
            "source": "default",
        }
        for question in INTAKE_QUESTIONS
    }
    if enhance_mode:
        # PPTX→HTML enhance mode: content is already finalised — push toward rich HTML output.
        for key, value, label in (
            ("output_format", "html-revealjs", "HTML（Reveal.js）"),
            ("visual_freedom", "delegate", "AI 自主决策"),
        ):
            default_selections[key] = {"value": value, "label": label, "source": "enhance-mode-default"}
    return {
        "version": "0.1",
        "task_slug": task_slug,
        "topic": topic or task_slug.replace("-", " "),
        "ui_language": resolved_ui_language,
        "ui_language_source": "explicit" if ui_language in SUPPORTED_UI_LANGUAGES else "auto-detected",
        "conversation_text": conversation_text,
        "sources": source_items,
        "selections": default_selections,
        "output_format": output_format_from_selections(default_selections, "html-revealjs" if enhance_mode else "pptx"),
        "risks": infer_risks(source_items),
        "confirmed": False,
        "enhance_mode": enhance_mode,
        "created_at": datetime.now().isoformat(timespec="seconds"),
    }


def infer_source_type(source: str) -> str:
    lowered: str = source.lower()
    if re.match(r"https?://", source) and (
        "drive.google.com" in lowered or "docs.google.com" in lowered
    ):
        return "google-drive-url"
    if re.match(r"https?://", source):
        return "url"
    path: Path = Path(source).expanduser()
    if path.is_dir():
        return "folder"
    suffix: str = path.suffix.lower().lstrip(".")
    return suffix or "text"


def build_source_items(sources: list[str]) -> list[JsonDict]:
    return [
        {"path": source, "priority": "primary", "type": infer_source_type(source)}
        for source in sources
        if source.strip()
    ]


def source_paths_from_items(sources: Any) -> list[str]:
    if not isinstance(sources, list):
        return []
    paths: list[str] = []
    for source in sources:
        if isinstance(source, dict):
            path: str = str(source.get("path", "")).strip()
            if path:
                paths.append(path)
    return paths


def parse_sources_text(value: str) -> list[str]:
    raw_items: list[str] = re.split(r"[\n,]+", value)
    seen: set[str] = set()
    sources: list[str] = []
    for raw_item in raw_items:
        source: str = raw_item.strip()
        if not source or source in seen:
            continue
        seen.add(source)
        sources.append(source)
    return sources


def format_sources_text(sources: Any) -> str:
    return "\n".join(source_paths_from_items(sources))


def infer_risks(sources: list[JsonDict]) -> list[str]:
    risks: list[str] = []
    source_text: str = " ".join(str(item.get("path", "")) for item in sources).lower()
    if "logo" not in source_text:
        risks.append("未发现明确 logo 文件；如需使用 logo，必须由用户提供或使用官方来源。")
    if not any(token in source_text for token in ("metric", "data", "csv", "xlsx", "指标", "数据")):
        risks.append("未发现明确量化数据文件；第一版可能需要用定性证明或标注缺失指标。")
    if not sources:
        risks.append("未提供具体资料路径；需要在生成前补充 source material。")
    return risks


def apply_intake_selection(draft: JsonDict, form: dict[str, list[str]]) -> JsonDict:
    selections: JsonDict = {}
    draft_selections: JsonDict = draft.get("selections", {})
    ui_language: str = ui_language_from_brief(draft)
    form_sources_text: str = first_form_value(form, "sources_text", "").strip()
    draft_sources: list[str] = source_paths_from_items(draft.get("sources", []))
    selected_sources: list[str] = parse_sources_text(form_sources_text) if form_sources_text else draft_sources
    selected_source_items: list[JsonDict] = build_source_items(selected_sources)
    for question in INTAKE_QUESTIONS:
        draft_item: Any = draft_selections.get(question.key, {})
        fallback: str = (
            str(draft_item.get("value", ""))
            if isinstance(draft_item, dict) and draft_item.get("value")
            else default_intake_value(question, selected_sources, ui_language)
        )
        value: str = first_form_value(form, question.key, fallback)
        choice: Choice = selected_choice(question, value)
        if value not in {item.value for item in question.choices}:
            value = choice.value
        custom: str = first_form_value(form, f"{question.key}__custom", "").strip()
        selections[question.key] = {
            "value": value,
            "label": custom if value == "custom" and custom else choice.label,
            "source": "user-selected",
            "description": choice.description,
        }
        if custom:
            selections[question.key]["custom"] = custom

    brief: JsonDict = dict(draft)
    brief["sources"] = selected_source_items
    brief["risks"] = infer_risks(selected_source_items)
    brief["selections"] = selections
    brief["output_format"] = output_format_from_selections(selections, "pptx")
    brief["updated_at"] = datetime.now().isoformat(timespec="seconds")
    brief["confirmed"] = False
    return brief


def apply_visual_selection(selected: JsonDict, form: dict[str, list[str]]) -> JsonDict:
    topic: str = str(selected.get("topic", ""))
    selections: JsonDict = selected.get("selections", {})
    candidates: tuple[VisualCandidate, ...] = build_visual_candidates(topic, selections)
    selected_key: str = first_form_value(form, "visual_candidate", candidates[0].key)
    candidate: VisualCandidate = next((item for item in candidates if item.key == selected_key), candidates[0])
    updated: JsonDict = dict(selected)
    updated["visual_direction"] = {
        "selected_candidate": visual_candidate_to_json(candidate),
        "available_candidates": [visual_candidate_to_json(item) for item in candidates],
        "source": "user-selected",
        "notes": first_form_value(form, "visual_notes", "").strip(),
        "selected_at": datetime.now().isoformat(timespec="seconds"),
    }
    updated["updated_at"] = datetime.now().isoformat(timespec="seconds")
    updated["confirmed"] = False
    return updated


def ensure_visual_selection(selected: JsonDict) -> JsonDict:
    visual_direction: Any = selected.get("visual_direction", {})
    if isinstance(visual_direction, dict) and isinstance(visual_direction.get("selected_candidate"), dict):
        return selected
    topic: str = str(selected.get("topic", ""))
    selections: JsonDict = selected.get("selections", {})
    candidate: VisualCandidate = build_visual_candidates(topic, selections)[0]
    updated: JsonDict = dict(selected)
    updated["visual_direction"] = {
        "selected_candidate": visual_candidate_to_json(candidate),
        "available_candidates": [visual_candidate_to_json(item) for item in build_visual_candidates(topic, selections)],
        "source": "agent-recommended-default",
        "notes": "",
        "selected_at": datetime.now().isoformat(timespec="seconds"),
    }
    return updated


def first_form_value(form: dict[str, list[str]], key: str, default: str) -> str:
    values: list[str] | None = form.get(key)
    if not values:
        return default
    return values[0]


def candidate_slide_dirs(version_dir: Path) -> list[Path]:
    return [
        version_dir / "slides",
        version_dir / "preview",
        version_dir / "previews",
        version_dir / "rendered-slides",
        version_dir / "html-assets",
    ]


def collect_slide_images(version_dir: Path, explicit_dir: Path | None = None) -> list[Path]:
    search_dirs: list[Path] = [explicit_dir] if explicit_dir else candidate_slide_dirs(version_dir)
    for slide_dir in search_dirs:
        if slide_dir is None or not slide_dir.exists() or not slide_dir.is_dir():
            continue
        images: list[Path] = [
            path
            for path in slide_dir.iterdir()
            if path.is_file() and path.suffix.lower() in IMAGE_EXTENSIONS
        ]
        if images:
            return sorted(images, key=natural_sort_key)
    return []


def image_data_uri(path: Path) -> str:
    mime_type: str = mimetypes.guess_type(path.name)[0] or "image/png"
    encoded: str = base64.b64encode(path.read_bytes()).decode("ascii")
    return f"data:{mime_type};base64,{encoded}"


def build_share_html(title: str, slide_images: list[Path], source_pptx: Path | None = None) -> str:
    generated_at: str = datetime.now().isoformat(timespec="seconds")
    slides: list[str] = []
    for index, image_path in enumerate(slide_images, start=1):
        slides.append(
            f"""<section class="slide" id="slide-{index}">
  <div class="slide-number">{index:02d} / {len(slide_images):02d}</div>
  <img src="{image_data_uri(image_path)}" alt="Slide {index}">
</section>"""
        )
    source_note: str = (
        f"<p>Source PPTX: <code>{html.escape(str(source_pptx))}</code></p>"
        if source_pptx
        else ""
    )
    return f"""<!doctype html>
<html lang="zh-CN">
<head>
  <meta charset="utf-8">
  <meta name="viewport" content="width=device-width, initial-scale=1">
  <title>{html.escape(title)}</title>
  <style>
    :root {{
      color-scheme: dark;
      --bg: #111;
      --panel: #1b1b1b;
      --ink: #f6f6f6;
      --muted: #aaa;
    }}
    * {{ box-sizing: border-box; }}
    body {{
      margin: 0;
      background: var(--bg);
      color: var(--ink);
      font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", "Noto Sans SC", sans-serif;
    }}
    header {{
      position: sticky;
      top: 0;
      z-index: 10;
      display: flex;
      align-items: center;
      justify-content: space-between;
      gap: 16px;
      padding: 12px 20px;
      background: rgba(17, 17, 17, .92);
      border-bottom: 1px solid #2a2a2a;
      backdrop-filter: blur(8px);
    }}
    h1 {{ margin: 0; font-size: 16px; font-weight: 650; }}
    .meta {{ color: var(--muted); font-size: 12px; }}
    main {{
      display: grid;
      gap: 28px;
      padding: 24px;
      max-width: 1280px;
      margin: 0 auto;
    }}
    .slide {{
      position: relative;
      background: var(--panel);
      border: 1px solid #2a2a2a;
      box-shadow: 0 18px 60px rgba(0, 0, 0, .35);
    }}
    .slide img {{
      display: block;
      width: 100%;
      height: auto;
    }}
    .slide-number {{
      position: absolute;
      right: 10px;
      top: 8px;
      padding: 3px 7px;
      border-radius: 999px;
      background: rgba(0, 0, 0, .56);
      color: #fff;
      font-size: 11px;
    }}
    footer {{
      padding: 22px 24px 36px;
      color: var(--muted);
      font-size: 12px;
      text-align: center;
    }}
    code {{ color: #ddd; }}
  </style>
</head>
<body>
  <header>
    <h1>{html.escape(title)}</h1>
    <div class="meta">{len(slide_images)} slides · generated {html.escape(generated_at)}</div>
  </header>
  <main>
    {''.join(slides)}
  </main>
  <footer>
    <p>View-only HTML companion generated from rendered slide previews. Edit the PPTX source, then regenerate this HTML after changes.</p>
    {source_note}
  </footer>
</body>
</html>
"""


def write_share_html(
    task_dir: Path,
    version_dir: Path,
    title: str,
    explicit_slides_dir: Path | None = None,
    output_path: Path | None = None,
) -> Path:
    slide_images: list[Path] = collect_slide_images(version_dir, explicit_slides_dir)
    if not slide_images:
        searched: str = ", ".join(str(path) for path in candidate_slide_dirs(version_dir))
        raise ValueError(f"No per-slide preview images found. Searched: {searched}")
    final_dir: Path = task_dir / "final"
    final_dir.mkdir(parents=True, exist_ok=True)
    target: Path = output_path or final_dir / f"{slugify(title)}.html"
    source_pptx: Path = version_dir / "final.pptx"
    html_text: str = build_share_html(title, slide_images, source_pptx if source_pptx.exists() else None)
    write_text(target, html_text)
    return target


def _render_pptx_structure_panel(structure: list[JsonDict], ui_language: str) -> str:
    if not structure:
        return ""
    slide_rows: str = "".join(
        f'<tr><td style="width:2.5em;text-align:center;color:#5a7080;font-size:0.85em">{s["slide_num"]}</td>'
        f'<td style="font-weight:600">{html.escape(str(s.get("title","—")))}</td>'
        f'<td style="color:#5a7080;font-size:0.85em">{html.escape(", ".join(s.get("bullets",[])[:3]))}</td></tr>'
        for s in structure
    )
    lang_labels: dict[str, str] = {
        "zh": "已从 PPTX 提取的幻灯片结构",
        "en": "Slide structure extracted from PPTX",
        "de": "Aus PPTX extrahierte Folienstruktur",
        "fr": "Structure extraite du PPTX",
        "it": "Struttura estratta dal PPTX",
        "es": "Estructura extraída del PPTX",
    }
    label: str = lang_labels.get(ui_language, lang_labels["en"])
    return (
        f'<section class="section" style="border-left:4px solid #0f6f8f;padding-left:1em;margin-bottom:1.5em">'
        f'<h2 style="color:#0f6f8f;margin-top:0">✦ {html.escape(label)}</h2>'
        f'<table style="width:100%;border-collapse:collapse;font-size:0.85em">'
        f'<thead><tr><th>#</th><th>Title</th><th>Content preview</th></tr></thead>'
        f'<tbody>{slide_rows}</tbody></table>'
        f'<p style="font-size:0.8em;color:#5a7080;margin-top:0.5em">'
        f'The generated HTML will be recreated from this structure with a richer visual style.</p>'
        f'</section>'
    )


def render_intake(task_dir: Path) -> str:
    draft: JsonDict = read_json(task_dir / "brief-draft.json")
    current: JsonDict = read_json(task_dir / "intake-selection.json", draft)
    ui_language: str = ui_language_from_brief(current)
    selections: JsonDict = current.get("selections", {})
    current_sources: Any = current.get("sources", draft.get("sources", []))
    source_list: str = render_sources(current_sources, ui_language)
    sources_text: str = format_sources_text(current_sources)
    enhance_mode: bool = bool(current.get("enhance_mode", draft.get("enhance_mode", False)))
    pptx_structure: list[JsonDict] = list(current.get("pptx_structure", draft.get("pptx_structure", [])))
    enhance_banner: str = ""
    if enhance_mode:
        lang_titles: dict[str, str] = {
            "zh": "增强模式：PPTX → 酷炫 HTML",
            "en": "Enhance Mode: PPTX → Rich HTML",
            "de": "Enhance-Modus: PPTX → Ansprechendes HTML",
            "fr": "Mode amélioration : PPTX → HTML enrichi",
            "it": "Modalità miglioramento: PPTX → HTML avanzato",
            "es": "Modo mejora: PPTX → HTML enriquecido",
        }
        lang_descs: dict[str, str] = {
            "zh": "内容已从 PPTX 提取完毕。选择视觉风格后，将生成带完整动画和渐变背景的 HTML 展示版本。",
            "en": "Content has been extracted from your PPTX. Choose a visual direction — the HTML will be regenerated with animations and gradient backgrounds.",
            "de": "Inhalt wurde aus der PPTX extrahiert. Wählen Sie eine visuelle Richtung für die animierte HTML-Version.",
            "fr": "Le contenu a été extrait de votre PPTX. Choisissez une direction visuelle pour la version HTML animée.",
            "it": "Il contenuto è stato estratto dal PPTX. Scegliere una direzione visiva per la versione HTML animata.",
            "es": "El contenido ha sido extraído del PPTX. Elija una dirección visual para la versión HTML animada.",
        }
        enhance_banner = (
            f'<div style="background:#f0f8fb;border:1px solid #0f6f8f;border-radius:6px;'
            f'padding:0.9em 1.2em;margin-bottom:1.5em">'
            f'<strong style="color:#0f6f8f">⚡ {html.escape(lang_titles.get(ui_language, lang_titles["en"]))}</strong><br>'
            f'<span style="font-size:0.9em;color:#132238">{html.escape(lang_descs.get(ui_language, lang_descs["en"]))}</span>'
            f'</div>'
        )
    structure_panel: str = _render_pptx_structure_panel(pptx_structure, ui_language) if enhance_mode else ""
    body: str = f"""{enhance_banner}<div class="topline">{html.escape(t(ui_language, "intake_topline"))}</div>
<h1>{html.escape(t(ui_language, "intake_title"))}</h1>
<p>{html.escape(t(ui_language, "intake_intro"))}</p>
{structure_panel}
<section class="section">
  <h2>{html.escape(t(ui_language, "source_material"))}</h2>
  {source_list}
  <label>{html.escape(t(ui_language, "source_paths_label"))}
    <textarea name="sources_text" form="intake-form" placeholder="{html.escape(t(ui_language, "source_paths_placeholder"))}">{html.escape(sources_text)}</textarea>
  </label>
  <p class="meta">{html.escape(t(ui_language, "source_paths_meta"))}</p>
  <label>{html.escape(t(ui_language, "topic_title_label"))}
    <input type="text" name="topic" form="intake-form" value="{html.escape(str(current.get("topic", draft.get("topic", ""))))}">
  </label>
</section>
<form method="post" action="/api/intake" id="intake-form">
  {''.join(question_section(question, selections.get(question.key, {}).get("value", ""), ui_language) for question in INTAKE_QUESTIONS)}
  <section class="section">
    <h2>{html.escape(t(ui_language, "extra_notes"))}</h2>
    <textarea name="notes" placeholder="{html.escape(t(ui_language, "extra_notes_placeholder"))}">{html.escape(str(current.get("notes", "")))}</textarea>
  </section>
  <div class="actions">
    <button type="submit">{html.escape(t(ui_language, "next_visual"))}</button>
  </div>
</form>"""
    return html_page(t(ui_language, "intake_title"), body, ui_language)


def render_visual_inspiration(task_dir: Path) -> str:
    draft: JsonDict = read_json(task_dir / "brief-draft.json")
    selected: JsonDict = read_json(task_dir / "intake-selection.json", draft)
    ui_language: str = ui_language_from_brief(selected)
    topic: str = str(selected.get("topic", draft.get("topic", "")))
    selections: JsonDict = selected.get("selections", {})
    candidates: tuple[VisualCandidate, ...] = build_visual_candidates(topic, selections)
    output_format: str = str(selected.get("output_format", output_format_from_selections(selections, "pptx")))
    show_html_fields: bool = output_format in {"html-revealjs", "both"}
    visual_direction: JsonDict = selected.get("visual_direction", {})
    current_key: str = str(
        visual_direction.get("selected_candidate", {}).get("key", candidates[0].key)
        if isinstance(visual_direction, dict)
        else candidates[0].key
    )
    candidate_cards: list[str] = [
        render_visual_candidate_card(candidate, current_key == candidate.key, ui_language, show_html_fields)
        for candidate in candidates
    ]
    body: str = f"""<div class="topline">{html.escape(t(ui_language, "visual_gate"))}</div>
<h1>{html.escape(t(ui_language, "visual_title"))}</h1>
<p>{html.escape(t(ui_language, "visual_intro"))}</p>
<section class="section">
  <h2>{html.escape(t(ui_language, "current_topic"))}</h2>
  <p><strong>{html.escape(topic)}</strong></p>
</section>
<form method="post" action="/api/visual-inspiration">
  <div class="grid">
    {''.join(candidate_cards)}
  </div>
  <section class="section">
    <h2>{html.escape(t(ui_language, "visual_notes"))}</h2>
    <textarea name="visual_notes" placeholder="{html.escape(t(ui_language, "visual_notes_placeholder"))}">{html.escape(str(visual_direction.get("notes", "") if isinstance(visual_direction, dict) else ""))}</textarea>
  </section>
  <div class="actions">
    <a class="button secondary" href="/intake">{html.escape(t(ui_language, "back_intake"))}</a>
    <button type="submit">{html.escape(t(ui_language, "next_confirm"))}</button>
  </div>
</form>"""
    return html_page(t(ui_language, "visual_title"), body, ui_language)


def render_visual_candidate_card(
    candidate: VisualCandidate,
    checked: bool,
    ui_language: str = "zh",
    show_html_fields: bool = False,
) -> str:
    is_checked: str = " checked" if checked else ""
    candidate_json: JsonDict = visual_candidate_to_json(candidate)
    swatches: str = "".join(
        f'<span class="swatch-preview" style="background:{html.escape(color)}"></span>'
        for color in candidate.palette
    )
    bg_color: str = candidate.palette[0]
    ink_color: str = candidate.palette[1]
    accent: str = candidate.palette[2]
    accent_2: str = candidate.palette[3]
    gradient_preview: str = ""
    if candidate.html_gradient:
        gradient_label: str = candidate.html_gradient[:30]
        gradient_preview = (
            f"""<div class="html-field">
    <span class="label">{html.escape(t(ui_language, "html_gradient"))}</span>
    <span class="value" style="background: {html.escape(candidate.html_gradient)}; color: white; padding: 2px 8px; border-radius: 4px;">{html.escape(gradient_label)}</span>
  </div>"""
        )
    html_fields: str = ""
    if show_html_fields:
        html_fields = f"""<div class="section" style="margin-top: 12px; padding: 10px;">
  <div class="html-field">
    <span class="label">{html.escape(t(ui_language, "html_transition"))}</span>
    <span class="value">{html.escape(candidate.html_transition)}</span>
  </div>
  <div class="html-field">
    <span class="label">{html.escape(t(ui_language, "html_animation"))}</span>
    <span class="value">{html.escape(candidate.html_animation)}</span>
  </div>
  {gradient_preview}
</div>"""
    return f"""<label class="option candidate">
  <input type="radio" name="visual_candidate" value="{html.escape(candidate.key)}"{is_checked}>
  <strong>{html.escape(candidate.name)}</strong>
  <p>{html.escape(localized_visual_field(candidate_json, "summary", ui_language))}</p>
  <div class="swatches-preview">{swatches}</div>
  <div class="mini-slides" aria-hidden="true">
    <div class="mini-slide" style="background:{html.escape(bg_color)}; color:{html.escape(ink_color)}">
      <strong>{html.escape(candidate.name)}</strong>
      <span style="background:{html.escape(accent)}"></span>
    </div>
    <div class="mini-slide" style="background:#fff; color:#1f2937">
      <strong>{html.escape(t(ui_language, "evidence_page"))}</strong>
      <div class="mini-bars">
        <i style="width:86%; background:{html.escape(accent)}"></i>
        <i style="width:62%; background:{html.escape(accent_2)}"></i>
        <i style="width:74%; background:{html.escape(ink_color)}"></i>
      </div>
    </div>
  </div>
  <p><strong>{html.escape(t(ui_language, "best_for"))}:</strong> {html.escape(localized_visual_field(candidate_json, "best_for", ui_language))}</p>
  <p><strong>{html.escape(t(ui_language, "background"))}:</strong> {html.escape(localized_visual_field(candidate_json, "background", ui_language))}</p>
  <p><strong>{html.escape(t(ui_language, "chart"))}:</strong> {html.escape(localized_visual_field(candidate_json, "chart", ui_language))}</p>
  {html_fields}
  <p><strong>{html.escape(t(ui_language, "inspiration"))}:</strong> {html.escape(localized_visual_field(candidate_json, "inspiration", ui_language))}</p>
  <p class="meta"><strong>{html.escape(t(ui_language, "risk"))}:</strong> {html.escape(localized_visual_field(candidate_json, "risk", ui_language))}</p>
</label>"""


def render_sources(sources: Any, ui_language: str = "zh") -> str:
    if not isinstance(sources, list) or not sources:
        return f"<p class='risk'>{html.escape(t(ui_language, 'no_sources'))}</p>"
    items: list[str] = []
    for source in sources:
        if not isinstance(source, dict):
            continue
        path: str = str(source.get("path", ""))
        source_type: str = str(source.get("type", "unknown"))
        items.append(f"<li><code>{html.escape(path)}</code> <span class='source-tag'>{html.escape(source_type)}</span></li>")
    return f"<ul>{''.join(items)}</ul>"


def render_confirm(task_dir: Path) -> str:
    draft: JsonDict = read_json(task_dir / "brief-draft.json")
    selected: JsonDict = read_json(task_dir / "intake-selection.json", draft)
    ui_language: str = ui_language_from_brief(selected)
    confirm_token: str = ensure_confirm_token(task_dir)
    rows: list[str] = []
    selections: JsonDict = selected.get("selections", {})
    output_format: str = str(selected.get("output_format", output_format_from_selections(selections, "pptx")))
    for question in INTAKE_QUESTIONS:
        raw_item: Any = selections.get(question.key, {})
        item: JsonDict
        if isinstance(raw_item, dict) and raw_item.get("label"):
            item = raw_item
        else:
            choice: Choice = selected_choice(question, question.default)
            item = {
                "value": choice.value,
                "label": choice.label,
                "source": "default",
            }
        rows.append(
            "<tr>"
            f"<th>{html.escape(localized_question_title(question, ui_language))}</th>"
            f"<td>{html.escape(localized_choice_label(question, item, ui_language))}</td>"
            f"<td><span class='source-tag'>{html.escape(localized_source(str(item.get('source', 'unknown')), ui_language))}</span></td>"
            "</tr>"
        )
    risks: list[str] = selected.get("risks", draft.get("risks", []))
    risk_html: str = "".join(f"<li>{html.escape(localized_risk(str(risk), ui_language))}</li>" for risk in risks) or f"<li>{html.escape(t(ui_language, 'no_risks'))}</li>"
    visual_direction: JsonDict = selected.get("visual_direction", {})
    selected_candidate: JsonDict = {}
    if isinstance(visual_direction, dict):
        raw_candidate: Any = visual_direction.get("selected_candidate", {})
        if isinstance(raw_candidate, dict):
            selected_candidate = raw_candidate
    if not selected_candidate:
        topic: str = str(selected.get("topic", draft.get("topic", "")))
        candidates: tuple[VisualCandidate, ...] = build_visual_candidates(topic, selections)
        selected_candidate = visual_candidate_to_json(candidates[0])
    palette_html: str = "".join(
        f'<span class="swatch-preview" style="background:{html.escape(str(color))}"></span>'
        for color in selected_candidate.get("palette", [])
    )
    body: str = f"""<div class="topline">{html.escape(t(ui_language, "brief_gate"))}</div>
<h1>{html.escape(t(ui_language, "confirm_title"))}</h1>
<p>{html.escape(t(ui_language, "confirm_intro"))}</p>
<section class="section">
  <h2>{html.escape(t(ui_language, "topic"))}</h2>
  <p><strong>{html.escape(str(selected.get("topic", draft.get("topic", ""))))}</strong></p>
</section>
<section class="section">
  <h2>{html.escape(t(ui_language, "sources"))}</h2>
  {render_sources(selected.get("sources", draft.get("sources", [])), ui_language)}
</section>
<section class="section">
  <h2>{html.escape(t(ui_language, "summary"))}</h2>
  <table>
    <thead><tr><th>{html.escape(t(ui_language, "item"))}</th><th>{html.escape(t(ui_language, "selection"))}</th><th>{html.escape(t(ui_language, "source"))}</th></tr></thead>
    <tbody>{''.join(rows)}</tbody>
  </table>
</section>
<section class="section">
  <h2>{html.escape(t(ui_language, "visual_direction"))}</h2>
  <h3>{html.escape(str(selected_candidate.get("name", "")))}</h3>
  <p>{html.escape(localized_visual_field(selected_candidate, "summary", ui_language))}</p>
  <div class="swatches-preview">{palette_html}</div>
  <p><strong>{html.escape(t(ui_language, "background"))}:</strong> {html.escape(localized_visual_field(selected_candidate, "background", ui_language))}</p>
  <p><strong>{html.escape(t(ui_language, "layout"))}:</strong> {html.escape(localized_visual_field(selected_candidate, "layout", ui_language))}</p>
  <p><strong>{html.escape(t(ui_language, "chart"))}:</strong> {html.escape(localized_visual_field(selected_candidate, "chart", ui_language))}</p>
  <p><strong>{html.escape(t(ui_language, "image_strategy"))}:</strong> {html.escape(localized_visual_field(selected_candidate, "image_strategy", ui_language))}</p>
  <p><strong>{html.escape(t(ui_language, "risk"))}:</strong> {html.escape(localized_visual_field(selected_candidate, "risk", ui_language))}</p>
</section>
<section class="section">
  <h2>{html.escape(t(ui_language, "pre_generation_risks"))}</h2>
  <ul>{risk_html}</ul>
</section>
<section class="section">
  <h2>{html.escape(t(ui_language, "generation_strategy"))}</h2>
  <p>{html.escape(generation_strategy_text(output_format, task_dir, ui_language))}</p>
</section>
<form method="post" action="/api/confirm">
  <input type="hidden" name="confirm_token" value="{html.escape(confirm_token)}">
  <div class="actions">
    <a class="button secondary" href="/visual-inspiration">{html.escape(t(ui_language, "back_visual"))}</a>
    <button type="submit">{html.escape(t(ui_language, "confirm_button"))}</button>
  </div>
</form>"""
    return html_page(t(ui_language, "confirm_title"), body, ui_language)


def render_style_review(task_dir: Path) -> str:
    ui_language: str = ui_language_for_task(task_dir)
    version_name: str = latest_review_version(task_dir)
    version_dir: Path = task_dir / version_name
    contact_sheet: Path = version_dir / "contact-sheet.png"
    qa_summary: Path = version_dir / "qa-summary.md"
    pptx_path: Path = version_dir / "final.pptx"
    image_html: str = (
        f'<img class="contact-sheet" src="/static/{html.escape(version_name)}/contact-sheet.png" alt="{html.escape(version_name)} contact sheet">'
        if contact_sheet.exists()
        else f"<p class='risk'>{html.escape(t(ui_language, 'missing_contact_sheet').format(path=f'{version_name}/contact-sheet.png'))}</p>"
    )
    qa_text: str = qa_summary.read_text(encoding="utf-8") if qa_summary.exists() else t(ui_language, "missing_qa_summary")
    body: str = f"""<div class="topline">{html.escape(t(ui_language, "style_review"))}</div>
<h1>{html.escape(t(ui_language, "style_title"))}</h1>
<p>{html.escape(t(ui_language, "style_intro").format(version_name=version_name))}</p>
<section class="section">
  <h2>{html.escape(t(ui_language, "current_version"))}</h2>
  {image_html}
  <p>{html.escape(t(ui_language, "pptx_label"))}: <code>{html.escape(str(pptx_path))}</code></p>
  <pre>{html.escape(qa_text[:2400])}</pre>
</section>
<form method="post" action="/api/revision">
  <input type="hidden" name="base_version" value="{html.escape(version_name)}">
  {''.join(question_section(group, ui_language=ui_language) for group in REVISION_GROUPS)}
  <section class="section">
    <h2>{html.escape(t(ui_language, "revision_count_title"))}</h2>
    <label class="option"><input type="radio" name="revision_count" value="0"> {html.escape(t(ui_language, "keep_current_version"))}</label>
    <label class="option"><input type="radio" name="revision_count" value="1" checked> {html.escape(t(ui_language, "one_revision"))}</label>
    <label class="option"><input type="radio" name="revision_count" value="2"> {html.escape(t(ui_language, "two_revisions"))}</label>
    <textarea name="notes" placeholder="{html.escape(t(ui_language, "revision_notes_placeholder"))}"></textarea>
  </section>
  <div class="actions">
    <button type="submit">{html.escape(t(ui_language, "confirm_visual_choice"))}</button>
  </div>
</form>"""
    return html_page(t(ui_language, "style_title"), body, ui_language)


def apply_revision_request(form: dict[str, list[str]]) -> JsonDict:
    base_version: str = first_form_value(form, "base_version", "v1").strip() or "v1"
    request: JsonDict = {
        "version": "0.1",
        "base_version": base_version,
        "revision_count": int(first_form_value(form, "revision_count", "1")),
        "created_at": datetime.now().isoformat(timespec="seconds"),
        "preserve": [
            "factual content",
            "slide claims",
            "source attribution",
            "official asset policy",
            "imagegen policy",
        ],
    }
    for group in REVISION_GROUPS:
        value: str = first_form_value(form, group.key, group.default)
        choice: Choice = selected_choice(group, value)
        custom: str = first_form_value(form, f"{group.key}__custom", "").strip()
        request[group.key] = custom if value == "custom" and custom else value
        request[f"{group.key}_label"] = custom if value == "custom" and custom else choice.label
    request["notes"] = first_form_value(form, "notes", "").strip()
    return request


def render_compare(task_dir: Path) -> str:
    ui_language: str = ui_language_for_task(task_dir)
    version_cards: list[str] = []
    for version in ("v1", "v2", "v3"):
        version_dir: Path = task_dir / version
        if not version_dir.exists():
            continue
        contact_sheet: Path = version_dir / "contact-sheet.png"
        qa_summary: Path = version_dir / "qa-summary.md"
        pptx_path: Path = version_dir / "final.pptx"
        image_html: str = (
            f'<img class="contact-sheet" src="/static/{version}/contact-sheet.png" alt="{version} contact sheet">'
            if contact_sheet.exists()
            else f"<p class='risk'>{html.escape(t(ui_language, 'no_contact_sheet'))}</p>"
        )
        qa_text: str = qa_summary.read_text(encoding="utf-8") if qa_summary.exists() else t(ui_language, "missing_qa_summary")
        version_cards.append(
            f"""<section class="section">
  <h2>{html.escape(version.upper())}</h2>
  {image_html}
  <p>{html.escape(t(ui_language, "pptx_label"))}: <code>{html.escape(str(pptx_path))}</code></p>
  <pre>{html.escape(qa_text[:1600])}</pre>
  <label class="option"><input type="radio" name="selected_version" value="{html.escape(version)}"> {html.escape(t(ui_language, "choose_version").format(version=version.upper()))}</label>
</section>"""
        )
    if not version_cards:
        version_cards.append(f"<p class='risk'>{html.escape(t(ui_language, 'no_versions'))}</p>")
    body: str = f"""<div class="topline">{html.escape(t(ui_language, "version_compare"))}</div>
<h1>{html.escape(t(ui_language, "compare_title"))}</h1>
<form method="post" action="/api/final-selection">
  {''.join(version_cards)}
  <section class="section">
    <h2>{html.escape(t(ui_language, "choose_after_action"))}</h2>
    <textarea name="notes" placeholder="{html.escape(t(ui_language, "final_notes_placeholder"))}"></textarea>
  </section>
  <div class="actions">
    <button type="submit">{html.escape(t(ui_language, "confirm_final_version"))}</button>
    <a class="button secondary" href="/style-review">{html.escape(t(ui_language, "continue_editing"))}</a>
  </div>
</form>"""
    return html_page(t(ui_language, "compare_title"), body, ui_language)


def render_all_pages(task_dir: Path) -> None:
    write_text(task_dir / "intake.html", render_intake(task_dir))
    write_text(task_dir / "visual-inspiration.html", render_visual_inspiration(task_dir))
    write_text(task_dir / "brief-confirm.html", render_confirm(task_dir))
    write_text(task_dir / "style-review.html", render_style_review(task_dir))
    write_text(task_dir / "compare.html", render_compare(task_dir))


def version_number(path: Path) -> int:
    name: str = path.name
    if len(name) > 1 and name[0] == "v" and name[1:].isdigit():
        return int(name[1:])
    return -1


def latest_review_version(task_dir: Path) -> str:
    candidates: list[Path] = [
        item
        for item in task_dir.iterdir()
        if item.is_dir()
        and version_number(item) >= 1
        and (item / "contact-sheet.png").exists()
        and (item / "final.pptx").exists()
    ]
    if not candidates:
        return "v1"
    return max(candidates, key=version_number).name


def initial_prompt(task_dir: Path) -> str:
    brief: JsonDict = read_json(task_dir / "brief-confirmed.json")
    if not brief:
        return "No confirmed brief found. Confirm intake first."
    script_path: Path = Path(__file__).resolve()
    output_format: str = str(brief.get("output_format", output_format_from_selections(brief.get("selections", {}), "pptx")))
    common_rules: str = f"""Confirmed brief:
{json.dumps(brief, ensure_ascii=False, indent=2)}

Rules:
- Before generating, run:
  python3 "{script_path}" --base-dir "{task_dir.parent.parent}" guard --task "{task_dir.name}"
  If the guard fails, open the confirmation page through serve-wait and continue only after the user's HTML click:
  python3 "{script_path}" --base-dir "{task_dir.parent.parent}" serve-wait --task "{task_dir.name}" --for confirmed
- Audience, goal, output_format, research strategy, source boundary, content language, logo policy, image policy, selected visual direction, and output constraints are locked.
- Do not fabricate metrics, logos, customer names, screenshots, or official-looking brand assets.
- Use official or user-provided brand assets only.
- AI images are allowed only according to the confirmed image policy.
- Composition, layout rhythm, chart treatment, typography hierarchy, and visual expression should follow the selected visual candidate.
- Do not use a fixed design-lock unless the confirmed brief explicitly asks for it.
"""

    html_output: Path = task_dir / "final" / f"{task_dir.name}.html"
    pptx_output: Path = task_dir / "v1" / "final.pptx"

    if output_format == "html-revealjs":
        return f"""Write a Reveal.js 5.1.0 HTML presentation directly. Do NOT call the Codex Presentations plugin.

{common_rules}

Reveal.js requirements:
- Use pinned CDN links for reveal.js@5.1.0 reset.css, reveal.css, a built-in theme, and reveal.js.
- Use `html_transition`, `html_animation`, and `html_gradient` from the selected visual candidate.
- Put speaker notes in `<aside class="notes">` on each slide.
- Keep the file browser-runnable and save it to {html_output}.
- Include PDF export guidance: append `?print-pdf`, print from Chrome/Edge, landscape, no headers/footers.

QA:
- Open in a browser or capture screenshots to verify slide navigation, gradients, speaker notes, and no text overflow.
- Return HTML path, screenshot/QA evidence, and remaining risks.
"""

    if output_format == "both":
        return f"""Generate both outputs: first editable PPTX via Codex Presentations, then Reveal.js 5.1.0 HTML directly.

{common_rules}

PPTX route:
- Use the Codex Presentations skill and artifact-tool presentation JSX.
- Use the Presentations internal scratch workspace as required by the plugin.
- Copy the editable PPTX to {pptx_output}.
- Copy per-slide preview PNGs to {task_dir / "v1" / "slides"}.
- Copy the contact sheet and a concise QA summary to {task_dir / "v1"}.
- Generate layout JSON and QA notes in the Presentations workspace.
- PPTX uses the same palette as the HTML direction, but with solid-color backgrounds instead of gradients.

HTML route:
- Write Reveal.js HTML directly; do NOT call Presentations plugin for HTML.
- Use pinned reveal.js@5.1.0 CDN links.
- Use `html_transition`, `html_animation`, and `html_gradient` from the selected visual candidate.
- Save the HTML deck to {html_output}.

QA:
- PPTX QA must include rendered no-overlap checks and a contact sheet.
- HTML QA must include browser load/navigation and text-overflow checks.
- Return PPTX path, HTML path, contact sheet path, QA summary, and remaining risks.
"""

    return f"""Use the Codex Presentations skill and artifact-tool presentation JSX.

{common_rules}

Output:
- Use the Presentations internal scratch workspace as required by the plugin.
- Copy the editable PPTX to {task_dir / "v1" / "final.pptx"}.
- Copy per-slide preview PNGs to {task_dir / "v1" / "slides"}.
- Copy the contact sheet and a concise QA summary to {task_dir / "v1"}.
- Generate layout JSON and QA notes in the Presentations workspace.
- Write a concise QA summary to {task_dir / "v1" / "qa-summary.md"}.
- QA must include a rendered no-overlap check: titles, subtitles, body text, labels, footers, page numbers, and connector lines must not collide.
- Prefer artifact-tool or headless rendering paths that do not trigger Microsoft PowerPoint file-access dialogs. If a PowerPoint-based render is unavoidable on macOS, start scripts/macos/powerpoint-grant-access-watcher.sh before rendering.
- Long titles must be checked after rendering; if a title wraps and covers the subtitle or body area, fix and re-render before handoff.
- After final selection, generate a view-only HTML companion at {task_dir / "final" / (task_dir.name + ".html")} from the selected version's per-slide previews.
- Before returning control to the user after v1 generation, regenerate Director pages and start the local Director server with the style-review page open:
  python3 "{script_path}" --base-dir "{task_dir.parent.parent}" serve --task "{task_dir.name}" --open-page style-review
- If waiting for a style decision, use serve-wait and continue from the revision/final-selection signal rather than asking the user to reply in chat.
- Return PPTX path, HTML companion path, contact sheet path, QA summary, and remaining risks.
"""


def revision_prompt(task_dir: Path) -> str:
    request: JsonDict = read_json(task_dir / "revision-request.json")
    if not request:
        return "No revision request found. Complete style review first."
    base_version: str = str(request.get("base_version") or "v1")
    base_pptx: Path = task_dir / base_version / "final.pptx"
    existing_versions: list[int] = [
        version_number(item)
        for item in task_dir.iterdir()
        if item.is_dir() and version_number(item) >= 1
    ]
    next_version_number: int = (max(existing_versions) if existing_versions else 1) + 1
    next_version: str = f"v{next_version_number}"
    return f"""Revise the existing {base_version} PPTX using the selected revision request.

Base version:
{base_pptx}

Revision request:
{json.dumps(request, ensure_ascii=False, indent=2)}

Preserve:
- factual content
- slide claims
- sources and omission notes
- official asset policy
- imagegen policy

Change only the selected visual directions:
- palette direction
- structure rhythm
- visual expression
- image/logo treatment

Render and QA:
- use the Presentations internal scratch workspace as required by the plugin
- copy revised PPTX under {task_dir / next_version / "final.pptx"} unless generating multiple versions
- copy per-slide preview PNGs into the version's `slides/` folder
- copy contact sheet and QA summary into the same version folder
- explicitly check rendered slides for text overlap, especially wrapped titles covering subtitles or body text
- fix any overlap/cropping/too-tight spacing and re-render affected slides before final selection
- compare against {base_version}
- regenerate the final view-only HTML companion from the selected version's per-slide previews after final selection
- document what changed and remaining risks
"""


class DirectorHandler(BaseHTTPRequestHandler):
    task_dir: Path

    def log_message(self, format_text: str, *args: Any) -> None:
        sys.stderr.write("[presentation-director] " + format_text % args + "\n")

    def do_GET(self) -> None:
        parsed = urlparse(self.path)
        path: str = parsed.path
        if path in ("/", "/intake"):
            self.send_html(render_intake(self.task_dir))
        elif path == "/visual-inspiration":
            self.send_html(render_visual_inspiration(self.task_dir))
        elif path == "/confirm":
            self.send_html(render_confirm(self.task_dir))
        elif path == "/style-review":
            self.send_html(render_style_review(self.task_dir))
        elif path == "/compare":
            self.send_html(render_compare(self.task_dir))
        elif path.startswith("/static/"):
            self.send_static(path.removeprefix("/static/"))
        elif path == "/confirmed":
            ui_language: str = ui_language_for_task(self.task_dir)
            self.send_html(message_page(t(ui_language, "confirmed_title"), t(ui_language, "confirmed_message"), ui_language))
        elif path == "/revision-saved":
            ui_language: str = ui_language_for_task(self.task_dir)
            self.send_html(message_page(t(ui_language, "revision_saved_title"), t(ui_language, "revision_saved_message"), ui_language))
        elif path == "/final-selected":
            ui_language: str = ui_language_for_task(self.task_dir)
            self.send_html(message_page(t(ui_language, "final_selected_title"), t(ui_language, "final_selected_message"), ui_language))
        else:
            self.send_error(HTTPStatus.NOT_FOUND)

    def do_POST(self) -> None:
        parsed = urlparse(self.path)
        form: dict[str, list[str]] = self.read_form()
        if parsed.path == "/api/intake":
            draft: JsonDict = read_json(self.task_dir / "brief-draft.json")
            brief: JsonDict = apply_intake_selection(draft, form)
            topic: str = first_form_value(form, "topic", "").strip()
            if topic:
                brief["topic"] = topic
            brief["notes"] = first_form_value(form, "notes", "").strip()
            write_json(self.task_dir / "intake-selection.json", brief)
            render_all_pages(self.task_dir)
            self.redirect("/visual-inspiration")
        elif parsed.path == "/api/visual-inspiration":
            selected: JsonDict = read_json(self.task_dir / "intake-selection.json")
            if not selected:
                selected = read_json(self.task_dir / "brief-draft.json")
            updated: JsonDict = apply_visual_selection(selected, form)
            write_json(self.task_dir / "intake-selection.json", updated)
            render_all_pages(self.task_dir)
            self.redirect("/confirm")
        elif parsed.path == "/api/confirm":
            confirm_token: str = first_form_value(form, "confirm_token", "")
            if not valid_confirm_token(self.task_dir, confirm_token):
                ui_language: str = ui_language_for_task(self.task_dir)
                self.send_error(
                    HTTPStatus.FORBIDDEN,
                    t(ui_language, "invalid_token"),
                )
                return
            selected: JsonDict = read_json(self.task_dir / "intake-selection.json")
            if not selected:
                selected = read_json(self.task_dir / "brief-draft.json")
            selected = ensure_visual_selection(selected)
            confirmed_at: str = datetime.now().isoformat(timespec="seconds")
            selected["confirmed"] = True
            selected["confirmed_at"] = confirmed_at
            selected["confirmation_gate"] = confirmation_receipt(confirm_token, confirmed_at)
            write_json(self.task_dir / "brief-confirmed.json", selected)
            touch_status(self.task_dir, "confirmed")
            render_all_pages(self.task_dir)
            self.redirect("/confirmed")
        elif parsed.path == "/api/revision":
            request: JsonDict = apply_revision_request(form)
            write_json(self.task_dir / "revision-request.json", request)
            touch_status(self.task_dir, "revision")
            render_all_pages(self.task_dir)
            self.redirect("/revision-saved")
        elif parsed.path == "/api/final-selection":
            selected_version: str = first_form_value(form, "selected_version", "v1")
            selected_version_dir: Path = self.task_dir / selected_version
            selected_pptx: Path = selected_version_dir / "final.pptx"
            final_dir: Path = self.task_dir / "final"
            final_pptx: Path = final_dir / f"{self.task_dir.name}.pptx"
            if selected_pptx.exists():
                final_dir.mkdir(parents=True, exist_ok=True)
                shutil.copy2(selected_pptx, final_pptx)
            final_html: Path | None = None
            share_html_error: str = ""
            try:
                final_html = write_share_html(self.task_dir, selected_version_dir, self.task_dir.name)
            except ValueError as exc:
                share_html_error = str(exc)
            payload: JsonDict = {
                "version": "0.1",
                "selected_version": selected_version,
                "selected_pptx": str(selected_pptx),
                "final_pptx": str(final_pptx if final_pptx.exists() else selected_pptx),
                "final_html": str(final_html) if final_html else "",
                "notes": first_form_value(form, "notes", "").strip(),
                "selected_at": datetime.now().isoformat(timespec="seconds"),
            }
            if share_html_error:
                payload["share_html_error"] = share_html_error
            write_json(self.task_dir / "final-selection.json", payload)
            touch_status(self.task_dir, "final-selection")
            self.redirect("/final-selected")
        else:
            self.send_error(HTTPStatus.NOT_FOUND)

    def read_form(self) -> dict[str, list[str]]:
        length: int = int(self.headers.get("Content-Length", "0"))
        body: bytes = self.rfile.read(length)
        return parse_qs(body.decode("utf-8"), keep_blank_values=True)

    def send_html(self, content: str) -> None:
        body: bytes = content.encode("utf-8")
        self.send_response(HTTPStatus.OK)
        self.send_header("Content-Type", "text/html; charset=utf-8")
        self.send_header("Content-Length", str(len(body)))
        self.end_headers()
        self.wfile.write(body)

    def send_static(self, raw_path: str) -> None:
        relative: Path = Path(unquote(raw_path))
        if relative.is_absolute() or ".." in relative.parts:
            self.send_error(HTTPStatus.BAD_REQUEST)
            return
        path: Path = self.task_dir / relative
        if not path.exists() or not path.is_file():
            self.send_error(HTTPStatus.NOT_FOUND)
            return
        content_type: str = "application/octet-stream"
        if path.suffix.lower() == ".png":
            content_type = "image/png"
        elif path.suffix.lower() == ".jpg" or path.suffix.lower() == ".jpeg":
            content_type = "image/jpeg"
        elif path.suffix.lower() == ".md":
            content_type = "text/plain; charset=utf-8"
        data: bytes = path.read_bytes()
        self.send_response(HTTPStatus.OK)
        self.send_header("Content-Type", content_type)
        self.send_header("Content-Length", str(len(data)))
        self.end_headers()
        self.wfile.write(data)

    def redirect(self, target: str) -> None:
        self.send_response(HTTPStatus.SEE_OTHER)
        self.send_header("Location", target)
        self.end_headers()


def message_page(title: str, message: str, ui_language: str = "zh") -> str:
    body: str = f"""<div class="topline">{html.escape(t(ui_language, "intake_topline"))}</div>
<h1>{html.escape(title)}</h1>
<section class="section"><p>{html.escape(message)}</p></section>
<div class="actions">
  <a class="button" href="/intake">{html.escape(t(ui_language, "nav_intake"))}</a>
  <a class="button" href="/visual-inspiration">{html.escape(t(ui_language, "nav_visual"))}</a>
  <a class="button" href="/style-review">{html.escape(t(ui_language, "nav_style"))}</a>
  <a class="button" href="/compare">{html.escape(t(ui_language, "nav_compare"))}</a>
</div>"""
    return html_page(title, body, ui_language)


def extract_pptx_structure(pptx_path: str) -> list[JsonDict]:
    """Extract slide titles and bullet text from a PPTX file using python-pptx.
    Returns an empty list if python-pptx is not installed or the file cannot be read.
    Caps at 5 bullets per slide to stay within content density limits.
    """
    try:
        from pptx import Presentation as _Presentation  # type: ignore
        from pptx.enum.text import PP_ALIGN  # noqa: F401  # optional; just test import
    except ImportError:
        return []
    try:
        prs = _Presentation(pptx_path)
    except Exception:
        return []
    slides: list[JsonDict] = []
    for i, slide in enumerate(prs.slides):
        title: str | None = None
        bullets: list[str] = []
        for shape in slide.shapes:
            if not hasattr(shape, "has_text_frame") or not shape.has_text_frame:
                continue
            text: str = shape.text_frame.text.strip()
            if not text:
                continue
            is_title: bool = False
            if hasattr(shape, "placeholder_format") and shape.placeholder_format is not None:
                ph_idx = getattr(shape.placeholder_format, "idx", None)
                if ph_idx == 0:
                    is_title = True
            if is_title:
                title = text
            elif not title:
                title = text
            else:
                for para in shape.text_frame.paragraphs:
                    pt = para.text.strip()
                    if pt and len(bullets) < 5:
                        bullets.append(pt)
        slides.append({
            "slide_num": i + 1,
            "title": title or f"Slide {i + 1}",
            "bullets": bullets,
        })
    return slides


def resolve_task_dir(args: argparse.Namespace) -> Path:
    base_dir: Path = Path(args.base_dir).expanduser().resolve()
    thread_id: str = args.thread_id or os.environ.get("CODEX_THREAD_ID") or now_id()
    task_slug: str = slugify(args.task)
    return workspace_root(base_dir, thread_id, task_slug)


def command_init(args: argparse.Namespace) -> None:
    task_dir: Path = resolve_task_dir(args)
    task_dir.mkdir(parents=True, exist_ok=True)
    enhance_mode: bool = getattr(args, "mode", "new") == "enhance"
    brief: JsonDict = build_draft_brief(
        slugify(args.task),
        args.topic or args.task,
        args.source or [],
        args.ui_language,
        args.conversation_text,
        enhance_mode=enhance_mode,
    )
    if enhance_mode:
        pptx_sources: list[str] = [s for s in (args.source or []) if s.lower().endswith(".pptx")]
        if pptx_sources:
            structure: list[JsonDict] = extract_pptx_structure(pptx_sources[0])
            if structure:
                brief["pptx_structure"] = structure
                print(f"Extracted {len(structure)} slides from {pptx_sources[0]}")
            else:
                print("Warning: could not extract PPTX structure (python-pptx missing or file unreadable).")
    write_json(task_dir / "brief-draft.json", brief)
    write_json(task_dir / "brief" / "draft-brief.json", brief)
    render_all_pages(task_dir)
    mode_label: str = " [ENHANCE MODE — PPTX→HTML]" if enhance_mode else ""
    print(f"Presentation Director task created{mode_label}: {task_dir}")
    print(f"Open intake page: {task_dir / 'intake.html'}")
    print("For click-to-submit flow, run:")
    print(f"  python3 scripts/presentation_director.py serve --task {slugify(args.task)}")


def command_render(args: argparse.Namespace) -> None:
    task_dir: Path = resolve_task_dir(args)
    render_all_pages(task_dir)
    print(f"Rendered pages in {task_dir}")
    if args.open_page:
        open_director_page(args.host, args.port, args.open_page)


def open_director_page(host: str, port: int, page: str) -> str:
    path: str | None = PAGE_PATHS.get(page)
    if path is None:
        raise SystemExit(f"Unknown page: {page}")
    url: str = f"http://{host}:{port}{path}"
    webbrowser.open(url)
    print(f"Opened {url}")
    return url


def command_serve(args: argparse.Namespace) -> None:
    task_dir: Path = resolve_task_dir(args)
    if not (task_dir / "brief-draft.json").exists():
        print(f"Missing brief-draft.json in {task_dir}. Run init first.", file=sys.stderr)
        raise SystemExit(1)
    handler_class: type[DirectorHandler] = type(
        "BoundDirectorHandler",
        (DirectorHandler,),
        {"task_dir": task_dir},
    )
    server: ThreadingHTTPServer = ThreadingHTTPServer((args.host, args.port), handler_class)
    print(f"Serving Presentation Director for {task_dir}")
    print(f"Intake:             http://{args.host}:{args.port}/intake")
    print(f"Visual inspiration: http://{args.host}:{args.port}/visual-inspiration")
    print(f"Confirm:            http://{args.host}:{args.port}/confirm")
    print(f"Style review:       http://{args.host}:{args.port}/style-review")
    print(f"Compare:            http://{args.host}:{args.port}/compare")
    if not args.no_open and args.open_page:
        open_director_page(args.host, args.port, args.open_page)
    try:
        server.serve_forever()
    except KeyboardInterrupt:
        print("\nStopping server.")
    finally:
        server.server_close()


def command_serve_wait(args: argparse.Namespace) -> None:
    task_dir: Path = resolve_task_dir(args)
    if not (task_dir / "brief-draft.json").exists():
        print(f"Missing brief-draft.json in {task_dir}. Run init first.", file=sys.stderr)
        raise SystemExit(1)

    filename: str | None = STATUS_FILES.get(args.for_status)
    if filename is None:
        raise SystemExit(f"Unknown status: {args.for_status}")
    target: Path = status_dir(task_dir) / filename
    if target.exists() and not args.allow_existing:
        target.unlink()

    handler_class: type[DirectorHandler] = type(
        "BoundDirectorHandler",
        (DirectorHandler,),
        {"task_dir": task_dir},
    )
    server: ThreadingHTTPServer = ThreadingHTTPServer((args.host, args.port), handler_class)
    thread = threading.Thread(target=server.serve_forever, kwargs={"poll_interval": 0.2}, daemon=True)
    thread.start()

    print(f"Serving Presentation Director for {task_dir}")
    print(f"Intake:             http://{args.host}:{args.port}/intake")
    print(f"Visual inspiration: http://{args.host}:{args.port}/visual-inspiration")
    print(f"Confirm:            http://{args.host}:{args.port}/confirm")
    print(f"Waiting for:        {target}")
    if not args.no_open and args.open_page:
        open_director_page(args.host, args.port, args.open_page)

    started: float = time.time()
    try:
        while True:
            if target.exists():
                print(f"Ready: {target}")
                return
            if args.timeout > 0 and time.time() - started > args.timeout:
                raise SystemExit(f"Timed out waiting for {target}")
            time.sleep(args.interval)
    except KeyboardInterrupt:
        print("\nStopping server.")
        raise
    finally:
        server.shutdown()
        server.server_close()
        thread.join(timeout=2)


def command_wait(args: argparse.Namespace) -> None:
    task_dir: Path = resolve_task_dir(args)
    filename: str | None = STATUS_FILES.get(args.for_status)
    if filename is None:
        raise SystemExit(f"Unknown status: {args.for_status}")
    target: Path = status_dir(task_dir) / filename
    started: float = time.time()
    while True:
        if target.exists():
            print(f"Ready: {target}")
            return
        if args.timeout > 0 and time.time() - started > args.timeout:
            raise SystemExit(f"Timed out waiting for {target}")
        time.sleep(args.interval)


def command_prompt(args: argparse.Namespace) -> None:
    task_dir: Path = resolve_task_dir(args)
    if args.kind == "initial":
        print(initial_prompt(task_dir))
    elif args.kind == "revision":
        print(revision_prompt(task_dir))
    else:
        raise SystemExit(f"Unknown prompt kind: {args.kind}")


def command_guard(args: argparse.Namespace) -> None:
    task_dir: Path = resolve_task_dir(args)
    errors: list[str] = validate_generation_guard(task_dir)
    if errors:
        print("Presentation Director guard failed:", file=sys.stderr)
        for error in errors:
            print(f"- {error}", file=sys.stderr)
        raise SystemExit(2)
    print(f"Presentation Director guard passed: {task_dir}")


def command_open_page(args: argparse.Namespace) -> None:
    resolve_task_dir(args)
    open_director_page(args.host, args.port, args.page)


def command_share_html(args: argparse.Namespace) -> None:
    task_dir: Path = resolve_task_dir(args)
    version_dir: Path = task_dir / args.version
    slides_dir: Path | None = Path(args.slides_dir).expanduser().resolve() if args.slides_dir else None
    output_path: Path | None = Path(args.output).expanduser().resolve() if args.output else None
    title: str = args.title or task_dir.name
    target: Path = write_share_html(task_dir, version_dir, title, slides_dir, output_path)
    print(f"Share HTML written: {target}")


def build_parser() -> argparse.ArgumentParser:
    parser: argparse.ArgumentParser = argparse.ArgumentParser(
        description="Presentation Director helper for Codex + Presentations workflows."
    )
    parser.add_argument("--base-dir", default=".", help="Repository/workspace root. Default: current directory.")
    parser.add_argument("--thread-id", default=None, help="Deprecated compatibility option; user-facing files now live in PPTX/<task-slug>/.")

    subparsers = parser.add_subparsers(dest="command", required=True)

    init_parser = subparsers.add_parser("init", help="Create a director workspace and render initial pages.")
    init_parser.add_argument("--task", required=True, help="Task slug or title.")
    init_parser.add_argument("--topic", default="", help="Optional topic/title shown in the brief.")
    init_parser.add_argument("--source", action="append", default=[], help="Source path or URL. Repeatable.")
    init_parser.add_argument(
        "--ui-language",
        choices=("auto", "zh", "en", "de", "fr", "it", "es"),
        default="auto",
        help="Language for Director communication UI. auto detects from --conversation-text or topic.",
    )
    init_parser.add_argument("--conversation-text", default="", help="Recent user conversation text for auto-detecting the Director UI language.")
    init_parser.add_argument(
        "--mode",
        choices=("new", "enhance"),
        default="new",
        help=(
            "Workflow mode. 'new' (default): full intake from source material. "
            "'enhance': extract content from an existing PPTX source and regenerate as a visually rich HTML presentation."
        ),
    )
    init_parser.set_defaults(func=command_init)

    render_parser = subparsers.add_parser("render", help="Regenerate HTML pages from current JSON.")
    render_parser.add_argument("--task", required=True, help="Task slug or title.")
    render_parser.add_argument("--host", default=DEFAULT_HOST, help=f"Host for optional browser open. Default: {DEFAULT_HOST}")
    render_parser.add_argument("--port", default=DEFAULT_PORT, type=int, help=f"Port for optional browser open. Default: {DEFAULT_PORT}")
    render_parser.add_argument("--open-page", choices=sorted(PAGE_PATHS.keys()), help="Open a Director page in the default browser after rendering.")
    render_parser.set_defaults(func=command_render)

    serve_parser = subparsers.add_parser("serve", help="Run local click-to-submit UI server.")
    serve_parser.add_argument("--task", required=True, help="Task slug or title.")
    serve_parser.add_argument("--host", default=DEFAULT_HOST, help=f"Host. Default: {DEFAULT_HOST}")
    serve_parser.add_argument("--port", default=DEFAULT_PORT, type=int, help=f"Port. Default: {DEFAULT_PORT}")
    serve_parser.add_argument(
        "--open-page",
        choices=sorted(PAGE_PATHS.keys()),
        default="intake",
        help="Open a Director page in the default browser once the server starts. Default: intake.",
    )
    serve_parser.add_argument("--no-open", action="store_true", help="Do not open a browser page after the server starts.")
    serve_parser.set_defaults(func=command_serve)

    serve_wait_parser = subparsers.add_parser(
        "serve-wait",
        help="Run the click-to-submit UI server, open a page, wait for a status signal, then stop the server.",
    )
    serve_wait_parser.add_argument("--task", required=True, help="Task slug or title.")
    serve_wait_parser.add_argument("--host", default=DEFAULT_HOST, help=f"Host. Default: {DEFAULT_HOST}")
    serve_wait_parser.add_argument("--port", default=DEFAULT_PORT, type=int, help=f"Port. Default: {DEFAULT_PORT}")
    serve_wait_parser.add_argument(
        "--open-page",
        choices=sorted(PAGE_PATHS.keys()),
        default="intake",
        help="Open a Director page in the default browser once the server starts. Default: intake.",
    )
    serve_wait_parser.add_argument("--no-open", action="store_true", help="Do not open a browser page after the server starts.")
    serve_wait_parser.add_argument("--for", dest="for_status", choices=sorted(STATUS_FILES.keys()), required=True)
    serve_wait_parser.add_argument("--timeout", type=float, default=0.0, help="Seconds before timeout. 0 means no timeout.")
    serve_wait_parser.add_argument("--interval", type=float, default=1.0, help="Polling interval seconds.")
    serve_wait_parser.add_argument(
        "--allow-existing",
        action="store_true",
        help="Treat an already-existing status file as ready instead of waiting for a fresh click.",
    )
    serve_wait_parser.set_defaults(func=command_serve_wait)

    wait_parser = subparsers.add_parser("wait", help="Wait for a ready status file.")
    wait_parser.add_argument("--task", required=True, help="Task slug or title.")
    wait_parser.add_argument("--for", dest="for_status", choices=sorted(STATUS_FILES.keys()), required=True)
    wait_parser.add_argument("--timeout", type=float, default=0.0, help="Seconds before timeout. 0 means no timeout.")
    wait_parser.add_argument("--interval", type=float, default=1.0, help="Polling interval seconds.")
    wait_parser.set_defaults(func=command_wait)

    prompt_parser = subparsers.add_parser("prompt", help="Print Presentations handoff prompt.")
    prompt_parser.add_argument("--task", required=True, help="Task slug or title.")
    prompt_parser.add_argument("--kind", choices=("initial", "revision"), required=True)
    prompt_parser.set_defaults(func=command_prompt)

    guard_parser = subparsers.add_parser("guard", help="Validate that a net-new PPTX task passed the user confirmation gate.")
    guard_parser.add_argument("--task", required=True, help="Task slug or title.")
    guard_parser.set_defaults(func=command_guard)

    open_parser = subparsers.add_parser("open-page", help="Open a running Director page in the default browser.")
    open_parser.add_argument("--task", required=True, help="Task slug or title.")
    open_parser.add_argument("--page", choices=sorted(PAGE_PATHS.keys()), required=True)
    open_parser.add_argument("--host", default=DEFAULT_HOST, help=f"Host. Default: {DEFAULT_HOST}")
    open_parser.add_argument("--port", default=DEFAULT_PORT, type=int, help=f"Port. Default: {DEFAULT_PORT}")
    open_parser.set_defaults(func=command_open_page)

    share_parser = subparsers.add_parser("share-html", help="Build a view-only final HTML companion from per-slide preview images.")
    share_parser.add_argument("--task", required=True, help="Task slug or title.")
    share_parser.add_argument("--version", default="v1", help="Version folder under PPTX/<task-slug>/. Default: v1.")
    share_parser.add_argument("--slides-dir", help="Optional explicit directory containing per-slide PNG/JPG/WebP previews.")
    share_parser.add_argument("--title", default="", help="Optional HTML title and output filename slug.")
    share_parser.add_argument("--output", help="Optional output HTML path. Default: PPTX/<task-slug>/final/<title>.html.")
    share_parser.set_defaults(func=command_share_html)

    return parser


def main() -> None:
    parser: argparse.ArgumentParser = build_parser()
    args: argparse.Namespace = parser.parse_args()
    args.func(args)


if __name__ == "__main__":
    main()
